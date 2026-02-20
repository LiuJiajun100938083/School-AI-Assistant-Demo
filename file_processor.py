# file_processor.py - 文件处理器模块
import os
import logging
from pathlib import Path
from typing import Tuple, Dict, Optional
import tempfile
import shutil

logger = logging.getLogger(__name__)


class FileProcessor:
    """文件处理器 - 处理各种格式的文档"""

    def __init__(self):
        self.supported_formats = {
            '.txt': self._process_text,
            '.md': self._process_text,
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.pptx': self._process_pptx
        }

        # 检查可选依赖
        self.pdf_available = self._check_pdf_support()
        self.docx_available = self._check_docx_support()
        self.pptx_available = self._check_pptx_support()

        logger.info("📄 文件处理器初始化完成")
        logger.info(f"  PDF支持: {'✅' if self.pdf_available else '❌'}")
        logger.info(f"  DOCX支持: {'✅' if self.docx_available else '❌'}")
        logger.info(f"  PPTX支持: {'✅' if self.pptx_available else '❌'}")

    def _check_pdf_support(self) -> bool:
        """检查PDF处理支持"""
        try:
            import PyPDF2
            return True
        except ImportError:
            try:
                import pdfplumber
                return True
            except ImportError:
                logger.warning("⚠️ PDF处理库未安装 (PyPDF2 或 pdfplumber)")
                return False

    def _check_docx_support(self) -> bool:
        """检查DOCX处理支持"""
        try:
            import docx
            return True
        except ImportError:
            logger.warning("⚠️ DOCX处理库未安装 (python-docx)")
            return False

    def _check_pptx_support(self) -> bool:
        """检查PPTX处理支持"""
        try:
            import pptx
            return True
        except ImportError:
            logger.warning("⚠️ PPTX处理库未安装 (python-pptx)")
            return False

    def is_supported_format(self, filename: str) -> bool:
        """检查是否支持的文件格式"""
        ext = Path(filename).suffix.lower()

        # 检查基本支持
        if ext not in self.supported_formats:
            return False

        # 检查特定格式的依赖
        if ext == '.pdf' and not self.pdf_available:
            return False
        if ext == '.docx' and not self.docx_available:
            return False
        if ext == '.pptx' and not self.pptx_available:
            return False

        return True

    def get_supported_formats(self) -> list:
        """获取支持的文件格式列表"""
        formats = []
        for ext in self.supported_formats:
            if ext in ['.txt', '.md']:
                formats.append(ext)
            elif ext == '.pdf' and self.pdf_available:
                formats.append(ext)
            elif ext == '.docx' and self.docx_available:
                formats.append(ext)
            elif ext == '.pptx' and self.pptx_available:
                formats.append(ext)
        return formats

    def process_file(self, file_path: str, original_filename: str,
                     user: str = "unknown") -> Tuple[bool, str, Dict]:
        """
        处理文件并提取文本内容

        Returns:
            (成功标志, 提取的文本内容, 文件信息字典)
        """
        try:
            # 获取文件扩展名
            ext = Path(original_filename).suffix.lower()

            # 检查是否支持
            if not self.is_supported_format(original_filename):
                return False, "", {"error": f"不支持的文件格式: {ext}"}

            # 调用对应的处理函数
            processor = self.supported_formats.get(ext)
            if processor:
                content = processor(file_path)

                # 获取文件信息
                file_stat = os.stat(file_path)
                file_info = {
                    "filename": original_filename,
                    "size": file_stat.st_size,
                    "format": ext,
                    "processed_by": user,
                    "content_length": len(content)
                }

                return True, content, file_info
            else:
                return False, "", {"error": f"没有找到处理器: {ext}"}

        except Exception as e:
            logger.error(f"❌ 处理文件失败 {original_filename}: {e}")
            return False, "", {"error": str(e)}

    def _process_text(self, file_path: str) -> str:
        """处理文本文件（TXT, MD）"""
        try:
            # 尝试不同的编码
            encodings = ['utf-8', 'gbk', 'gb2312', 'big5', 'utf-16']

            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                        logger.info(f"✅ 使用 {encoding} 编码成功读取文本文件")
                        return content
                except UnicodeDecodeError:
                    continue

            # 如果所有编码都失败，使用错误处理
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
                logger.warning("⚠️ 使用UTF-8编码（忽略错误）读取文本文件")
                return content

        except Exception as e:
            logger.error(f"❌ 读取文本文件失败: {e}")
            raise

    def _process_pdf(self, file_path: str) -> str:
        """处理PDF文件 - 增强版，支持多种方法和OCR"""
        if not self.pdf_available:
            raise Exception("PDF处理库未安装")

        content = ""

        # 方法1：尝试使用pdfplumber（更好的表格支持）
        try:
            import pdfplumber
            logger.info("🔄 尝试使用pdfplumber提取PDF...")
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            content += f"[第{i + 1}页]\n{page_text}\n\n"
                    except Exception as e:
                        logger.warning(f"⚠️ 第{i + 1}页提取失败: {e}")

            if content.strip():
                logger.info(f"✅ pdfplumber成功提取 {len(content)} 字符")
                return self._clean_pdf_text(content)
        except ImportError:
            logger.warning("⚠️ pdfplumber未安装")
        except Exception as e:
            logger.warning(f"⚠️ pdfplumber处理失败: {e}")

        # 方法2：尝试使用PyPDF2
        try:
            import PyPDF2
            logger.info("🔄 尝试使用PyPDF2提取PDF...")
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                num_pages = len(reader.pages)
                logger.info(f"📄 PDF共 {num_pages} 页")

                for page_num in range(num_pages):
                    try:
                        page = reader.pages[page_num]
                        page_text = page.extract_text()
                        if page_text:
                            content += f"[第{page_num + 1}页]\n{page_text}\n\n"
                    except Exception as e:
                        logger.warning(f"⚠️ 第{page_num + 1}页提取失败: {e}")

            if content.strip():
                logger.info(f"✅ PyPDF2成功提取 {len(content)} 字符")
                return self._clean_pdf_text(content)
        except ImportError:
            logger.warning("⚠️ PyPDF2未安装")
        except Exception as e:
            logger.warning(f"⚠️ PyPDF2处理失败: {e}")

        # 方法3：尝试使用pymupdf (fitz)
        try:
            import fitz  # PyMuPDF
            logger.info("🔄 尝试使用PyMuPDF提取PDF...")
            pdf_document = fitz.open(file_path)

            for page_num in range(pdf_document.page_count):
                try:
                    page = pdf_document[page_num]
                    page_text = page.get_text()
                    if page_text:
                        content += f"[第{page_num + 1}页]\n{page_text}\n\n"
                except Exception as e:
                    logger.warning(f"⚠️ 第{page_num + 1}页提取失败: {e}")

            pdf_document.close()

            if content.strip():
                logger.info(f"✅ PyMuPDF成功提取 {len(content)} 字符")
                return self._clean_pdf_text(content)
        except ImportError:
            logger.warning("⚠️ PyMuPDF未安装")
        except Exception as e:
            logger.warning(f"⚠️ PyMuPDF处理失败: {e}")

        # 方法4：如果以上都失败，尝试OCR（需要额外安装）
        if not content.strip():
            logger.warning("⚠️ 常规方法无法提取文本，尝试OCR...")
            content = self._try_ocr_pdf(file_path)
            if content:
                logger.info(f"✅ OCR成功提取 {len(content)} 字符")
                return content

        # 如果还是没有内容，但是PDF确实存在且可以打开
        try:
            # 检查文件是否真的是PDF
            with open(file_path, 'rb') as f:
                header = f.read(4)
                if header == b'%PDF':
                    # 是有效的PDF，但可能是扫描版或加密的
                    error_msg = """
这个PDF文件可能是：
1. 扫描版PDF（图片形式）- 需要OCR识别
2. 加密或受保护的PDF
3. 使用了特殊编码的PDF

解决方案：
1. 安装OCR支持: pip install pytesseract pdf2image
2. 转换PDF为文本格式后再上传
3. 使用在线PDF转文本工具处理后上传

暂时建议：请将内容复制为TXT或MD格式后上传。
                    """
                    logger.warning(error_msg)
                    # 返回一个占位符，避免完全失败
                    return f"[PDF文件: {Path(file_path).name}]\n[提取失败 - 可能是扫描版PDF]\n[建议转换为文本格式后重新上传]"
        except:
            pass

        raise Exception("PDF文件无法提取文本内容（可能是扫描版或加密PDF）")

    def _clean_pdf_text(self, text: str) -> str:
        """清理PDF提取的文本"""
        # 移除多余的空白字符
        lines = text.split('\n')
        cleaned_lines = []

        for line in lines:
            line = line.strip()
            if line:
                # 处理中文PDF的常见问题
                # 移除零宽字符
                line = line.replace('\u200b', '').replace('\ufeff', '')
                # 规范化空格
                line = ' '.join(line.split())
                cleaned_lines.append(line)

        # 重新组合，保持段落结构
        result = '\n'.join(cleaned_lines)

        # 处理常见的乱码模式
        result = result.replace('？', '?')  # 中文问号
        result = result.replace('！', '!')  # 中文感叹号

        return result

    def _try_ocr_pdf(self, file_path: str) -> str:
        """尝试使用OCR识别PDF（需要额外依赖）"""
        try:
            # 检查是否有OCR支持
            import pytesseract
            from pdf2image import convert_from_path

            logger.info("🔄 开始OCR识别...")

            # 将PDF转换为图片
            images = convert_from_path(file_path)

            content = []
            for i, image in enumerate(images):
                logger.info(f"🔍 OCR识别第 {i + 1}/{len(images)} 页...")
                # 使用OCR识别，支持中文
                text = pytesseract.image_to_string(image, lang='chi_sim+chi_tra+eng')
                if text.strip():
                    content.append(f"[第{i + 1}页]\n{text}")

            return '\n\n'.join(content)

        except ImportError:
            logger.warning("""
⚠️ OCR功能未安装。如需处理扫描版PDF，请安装：
1. 安装tesseract: 
   - Ubuntu/Debian: sudo apt-get install tesseract-ocr tesseract-ocr-chi-sim tesseract-ocr-chi-tra
   - macOS: brew install tesseract tesseract-lang
   - Windows: 下载安装 https://github.com/UB-Mannheim/tesseract/wiki
2. 安装Python库: pip install pytesseract pdf2image pillow
            """)
            return ""
        except Exception as e:
            logger.warning(f"⚠️ OCR识别失败: {e}")
            return ""

    def _process_docx(self, file_path: str) -> str:
        """处理DOCX文件"""
        if not self.docx_available:
            raise Exception("DOCX处理库未安装")

        try:
            import docx
            doc = docx.Document(file_path)

            content = []

            # 提取段落
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content.append(paragraph.text)

            # 提取表格
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        table_text.append(' | '.join(row_text))
                if table_text:
                    content.append('\n'.join(table_text))

            result = '\n\n'.join(content)
            logger.info(f"✅ 成功提取DOCX内容")
            return result

        except ImportError:
            raise Exception("请安装python-docx库: pip install python-docx")
        except Exception as e:
            logger.error(f"❌ DOCX处理失败: {e}")
            raise

    def _process_pptx(self, file_path: str) -> str:
        """处理PPTX文件"""
        if not self.pptx_available:
            raise Exception("PPTX处理库未安装")

        try:
            import pptx
            prs = pptx.Presentation(file_path)

            content = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = [f"=== 幻灯片 {slide_num} ==="]

                # 提取文本框内容
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)

                if len(slide_content) > 1:  # 如果有内容（不只是标题）
                    content.append('\n'.join(slide_content))

            result = '\n\n'.join(content)
            logger.info(f"✅ 成功提取PPTX内容")
            return result

        except ImportError:
            raise Exception("请安装python-pptx库: pip install python-pptx")
        except Exception as e:
            logger.error(f"❌ PPTX处理失败: {e}")
            raise

    def process_uploaded_file(self, file_data: bytes, filename: str,
                              user: str = "unknown") -> Tuple[bool, str, Dict]:
        """
        处理上传的文件数据

        Args:
            file_data: 文件二进制数据
            filename: 原始文件名
            user: 上传用户

        Returns:
            (成功标志, 提取的文本内容, 文件信息字典)
        """
        # 创建临时文件
        temp_file = None
        try:
            # 获取文件扩展名
            ext = Path(filename).suffix.lower()

            # 创建临时文件
            with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
                tmp.write(file_data)
                temp_file = tmp.name

            # 处理文件
            return self.process_file(temp_file, filename, user)

        finally:
            # 清理临时文件
            if temp_file and os.path.exists(temp_file):
                try:
                    os.unlink(temp_file)
                except:
                    pass


# 创建全局文件处理器实例
file_processor = FileProcessor()

# 测试代码
if __name__ == "__main__":
    print("🧪 测试文件处理器...")

    # 显示支持的格式
    supported = file_processor.get_supported_formats()
    print(f"📄 支持的格式: {', '.join(supported)}")

    # 测试文本文件
    test_file = "test.txt"
    if os.path.exists(test_file):
        success, content, info = file_processor.process_file(test_file, test_file)
        if success:
            print(f"✅ 成功处理 {test_file}")
            print(f"   内容长度: {info['content_length']} 字符")
        else:
            print(f"❌ 处理失败: {info.get('error', '未知错误')}")
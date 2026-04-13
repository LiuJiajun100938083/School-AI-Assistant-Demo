#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT 处理引擎

负责 PPT 文件的验证、转换和文本提取:
1. 文件安全验证 (类型、大小、完整性)
2. 使用 LibreOffice 将每页转换为 PNG 图片
3. 使用 python-pptx 提取每页文字内容 (供 AI 上下文)
4. 生成缩略图
"""

import asyncio
import logging
import os
import shutil
import subprocess
import uuid
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from PIL import Image

logger = logging.getLogger(__name__)

# PPT 文件头 (PK zip format)
_PPTX_MAGIC_BYTES = b"PK"

# 支持的 MIME 类型
_ALLOWED_MIME_TYPES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
}

# LibreOffice 可执行文件名称候选列表
_LIBREOFFICE_NAMES = ["soffice", "libreoffice"]

# Mac 上的 LibreOffice 典型路径
_LIBREOFFICE_MAC_PATHS = [
    "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    "/usr/local/bin/soffice",
    "/usr/local/bin/libreoffice",
]


def _find_libreoffice() -> Optional[str]:
    """自动检测 LibreOffice 路径"""
    # 优先检查 PATH 中的命令
    for name in _LIBREOFFICE_NAMES:
        path = shutil.which(name)
        if path:
            return path

    # 检查 Mac 特定路径
    for path in _LIBREOFFICE_MAC_PATHS:
        if os.path.isfile(path) and os.access(path, os.X_OK):
            return path

    return None


class PPTProcessor:
    """
    PPT 处理引擎

    用法:
        processor = PPTProcessor(upload_dir="./uploads/ppt", max_size_mb=150)
        processor.validate_file(file_bytes, "lesson.pptx")
        result = await processor.process_ppt(file_id, stored_path)
    """

    def __init__(
        self,
        upload_dir: str = "./uploads/ppt",
        max_size_mb: int = 150,
        image_width: int = 1920,
        thumbnail_width: int = 320,
    ):
        self.upload_dir = Path(upload_dir)
        self.max_size_bytes = max_size_mb * 1024 * 1024
        self.image_width = image_width
        self.thumbnail_width = thumbnail_width
        self.libreoffice_path = _find_libreoffice()

        if self.libreoffice_path:
            logger.info("LibreOffice 路径: %s", self.libreoffice_path)
        else:
            logger.warning(
                "未检测到 LibreOffice，PPT 将使用 python-pptx 降级处理"
            )

        # 确保目录存在
        self.upload_dir.mkdir(parents=True, exist_ok=True)

    # ============================================================
    # 文件验证
    # ============================================================

    def validate_file(
        self,
        file_bytes: bytes,
        filename: str,
    ) -> Tuple[bool, str]:
        """
        验证上传文件的安全性

        Returns:
            (is_valid, error_message)
        """
        # 1. 检查文件扩展名
        ext = Path(filename).suffix.lower()
        if ext not in (".pptx", ".ppt"):
            return False, f"不支持的文件类型: {ext}，仅支持 .pptx"

        # 2. 检查文件大小
        if len(file_bytes) > self.max_size_bytes:
            size_mb = len(file_bytes) / (1024 * 1024)
            max_mb = self.max_size_bytes / (1024 * 1024)
            return False, f"文件大小 {size_mb:.1f}MB 超出限制 {max_mb:.0f}MB"

        # 3. 检查 magic bytes (PK zip header)
        if not file_bytes[:2] == _PPTX_MAGIC_BYTES:
            return False, "文件格式无效 (非有效的 PPTX 文件)"

        # 4. 尝试用 python-pptx 打开验证完整性
        try:
            import io
            prs = Presentation(io.BytesIO(file_bytes))
            slide_count = len(prs.slides)
            if slide_count == 0:
                return False, "PPT 文件中没有幻灯片"
        except Exception as e:
            return False, f"PPT 文件损坏或格式不正确: {e}"

        return True, ""

    # ============================================================
    # 文件存储
    # ============================================================

    def save_file(
        self,
        file_bytes: bytes,
        file_id: str,
        original_filename: str,
    ) -> str:
        """
        安全保存上传文件

        文件存放位置: {upload_dir}/{file_id}/{file_id}.pptx
        使用 UUID 重命名，防止路径遍历攻击。

        Returns:
            存储路径 (相对路径)
        """
        file_dir = self.upload_dir / file_id
        file_dir.mkdir(parents=True, exist_ok=True)

        # 使用 file_id 作为文件名，保留原始扩展名
        ext = Path(original_filename).suffix.lower()
        stored_name = f"{file_id}{ext}"
        stored_path = file_dir / stored_name

        with open(stored_path, "wb") as f:
            f.write(file_bytes)

        logger.info(
            "PPT 文件已保存: %s (%d bytes)",
            stored_path, len(file_bytes),
        )
        return str(stored_path)

    # ============================================================
    # PPT 处理 (转图片 + 提取文字)
    # ============================================================

    async def process_ppt(
        self,
        file_id: str,
        stored_path: str,
    ) -> Dict[str, Any]:
        """
        处理 PPT 文件 (后台任务)

        步骤:
        1. 用 LibreOffice 将 PPT 每页转为 PNG
        2. 用 python-pptx 提取每页文字
        3. 生成缩略图
        4. 返回页面数据列表

        Returns:
            {
                "file_id": str,
                "total_pages": int,
                "pages": [
                    {
                        "page_id": str,
                        "page_number": int,
                        "image_path": str,
                        "thumbnail_path": str,
                        "text_content": str,
                    }, ...
                ]
            }
        """
        output_dir = self.upload_dir / file_id
        output_dir.mkdir(parents=True, exist_ok=True)

        # 提取文字 (python-pptx)
        texts = self._extract_all_texts(stored_path)

        # 转换为图片
        if self.libreoffice_path:
            image_paths = await self._convert_with_libreoffice(
                stored_path, output_dir
            )
        else:
            image_paths = self._convert_with_pillow_fallback(
                stored_path, output_dir
            )

        total_pages = max(len(image_paths), len(texts))

        # 构建页面数据
        pages = []
        for i in range(total_pages):
            page_id = str(uuid.uuid4())
            page_number = i + 1

            # 图片路径
            if i < len(image_paths):
                image_path = image_paths[i]
            else:
                image_path = ""

            # 生成缩略图
            thumbnail_path = ""
            if image_path and os.path.isfile(image_path):
                thumbnail_path = self._create_thumbnail(
                    image_path, output_dir, page_number
                )

            # 文字内容
            text_content = texts[i] if i < len(texts) else ""

            pages.append({
                "page_id": page_id,
                "page_number": page_number,
                "image_path": image_path,
                "thumbnail_path": thumbnail_path,
                "text_content": text_content,
            })

        logger.info(
            "PPT 处理完成: %s, %d 页",
            file_id, total_pages,
        )

        return {
            "file_id": file_id,
            "total_pages": total_pages,
            "pages": pages,
        }

    # ============================================================
    # LibreOffice 转换
    # ============================================================

    async def _convert_with_libreoffice(
        self,
        pptx_path: str,
        output_dir: Path,
    ) -> List[str]:
        """用 LibreOffice 将 PPT 每页转为 PNG（走 PDF 中间格式）"""
        try:
            return await self._convert_via_pdf(pptx_path, output_dir)
        except asyncio.TimeoutError:
            logger.error("LibreOffice 转换超时 (120s)")
            return self._convert_with_pillow_fallback(pptx_path, output_dir)
        except Exception as e:
            logger.error("LibreOffice 转换异常: %s", e)
            return self._convert_with_pillow_fallback(pptx_path, output_dir)

    async def _convert_via_pdf(
        self,
        pptx_path: str,
        output_dir: Path,
    ) -> List[str]:
        """
        LibreOffice: PPTX → PDF → 每页 PNG

        这是更可靠的多页转换方式。
        """
        # Step 1: PPTX → PDF
        # 每次转换用独立用户配置目录，避免并发锁冲突
        user_install = output_dir / ".lo_profile"
        cmd_pdf = [
            self.libreoffice_path,
            "--headless",
            "--norestore",
            "--nolockcheck",
            f"-env:UserInstallation=file://{user_install}",
            "--convert-to", "pdf",
            "--outdir", str(output_dir),
            pptx_path,
        ]

        process = await asyncio.create_subprocess_exec(
            *cmd_pdf,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        stdout, stderr = await asyncio.wait_for(
            process.communicate(), timeout=120,
        )

        # 清理临时配置目录
        shutil.rmtree(user_install, ignore_errors=True)

        if process.returncode != 0:
            logger.error(
                "LibreOffice PPTX→PDF 失败 (rc=%d): %s",
                process.returncode,
                stderr.decode("utf-8", errors="replace")[:500],
            )

        # 找到生成的 PDF 文件
        pdf_name = Path(pptx_path).stem + ".pdf"
        pdf_path = output_dir / pdf_name

        if not pdf_path.exists():
            logger.warning("PDF 未生成，降级处理")
            return self._convert_with_pillow_fallback(pptx_path, output_dir)

        # Step 2: PDF → 每页 PNG (使用 PyMuPDF)
        image_paths = []
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(str(pdf_path))
            for page_idx in range(len(doc)):
                page = doc[page_idx]
                # 高分辨率渲染
                zoom = self.image_width / page.rect.width
                mat = fitz.Matrix(zoom, zoom)
                pix = page.get_pixmap(matrix=mat)

                img_filename = f"page_{page_idx + 1}.png"
                img_path = str(output_dir / img_filename)
                pix.save(img_path)
                image_paths.append(img_path)

            doc.close()
            logger.info("PDF → PNG 转换完成: %d 页", len(image_paths))

        except ImportError:
            logger.warning("PyMuPDF 未安装，尝试 pdf2image")
            try:
                from pdf2image import convert_from_path
                images = convert_from_path(
                    str(pdf_path),
                    dpi=200,
                    fmt="png",
                )
                for idx, img in enumerate(images):
                    img_filename = f"page_{idx + 1}.png"
                    img_path = str(output_dir / img_filename)
                    img.save(img_path, "PNG")
                    image_paths.append(img_path)
            except Exception as e:
                logger.error("pdf2image 转换失败: %s", e)
                return self._convert_with_pillow_fallback(pptx_path, output_dir)

        except Exception as e:
            logger.error("PDF → PNG 失败: %s", e)
            return self._convert_with_pillow_fallback(pptx_path, output_dir)
        finally:
            # 清理中间 PDF
            try:
                pdf_path.unlink(missing_ok=True)
            except Exception:
                pass

        return image_paths

    # ============================================================
    # 降级方案: python-pptx + Pillow
    # ============================================================

    def _convert_with_pillow_fallback(
        self,
        pptx_path: str,
        output_dir: Path,
    ) -> List[str]:
        """
        降级方案: 无 LibreOffice 时，用纯 Python 生成简易页面图片

        注意: 这只能渲染文字内容，复杂排版/图片/动画会丢失。
        """
        logger.info("使用降级方案生成 PPT 页面图片")
        image_paths = []

        try:
            prs = Presentation(pptx_path)
            slide_width = prs.slide_width.pt if prs.slide_width else 960
            slide_height = prs.slide_height.pt if prs.slide_height else 540

            for idx, slide in enumerate(prs.slides):
                img_filename = f"page_{idx + 1}.png"
                img_path = str(output_dir / img_filename)

                # 创建白色背景图片
                img = Image.new(
                    "RGB",
                    (self.image_width, int(self.image_width * slide_height / slide_width)),
                    color=(255, 255, 255),
                )

                # 尝试在图片上写文字
                try:
                    from PIL import ImageDraw, ImageFont

                    draw = ImageDraw.Draw(img)
                    text = self._extract_slide_text(slide)

                    # 尝试加载系统字体 (优先 CJK 字体以支持中文)
                    font = None
                    for font_path, font_size in [
                        # macOS
                        ("/System/Library/Fonts/PingFang.ttc", 32),
                        # Linux: Noto CJK (Docker apt: fonts-noto-cjk)
                        ("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc", 30),
                        ("/usr/share/fonts/noto-cjk/NotoSansCJK-Regular.ttc", 30),
                        # Linux: WenQuanYi
                        ("/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc", 28),
                        # Fallback Latin
                        ("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 28),
                    ]:
                        try:
                            font = ImageFont.truetype(font_path, font_size)
                            break
                        except Exception:
                            continue
                    if font is None:
                        font = ImageFont.load_default()

                    # 绘制标题
                    title_text = f"Slide {idx + 1}"
                    if slide.shapes.title and slide.shapes.title.text:
                        title_text = slide.shapes.title.text

                    draw.text((60, 40), title_text, fill=(0, 0, 0), font=font)

                    # 绘制正文
                    if text:
                        lines = text.split("\n")[:20]
                        y_offset = 120
                        for line in lines:
                            if line.strip():
                                draw.text(
                                    (60, y_offset),
                                    line[:80],
                                    fill=(51, 51, 51),
                                    font=font,
                                )
                                y_offset += 45
                except Exception as e:
                    logger.debug("绘制文字失败: %s", e)

                img.save(img_path, "PNG")
                image_paths.append(img_path)

        except Exception as e:
            logger.error("降级方案处理失败: %s", e)

        return image_paths

    # ============================================================
    # 文本提取
    # ============================================================

    def _extract_all_texts(self, pptx_path: str) -> List[str]:
        """提取所有幻灯片的文字内容"""
        texts = []
        try:
            prs = Presentation(pptx_path)
            for slide in prs.slides:
                text = self._extract_slide_text(slide)
                texts.append(text)
        except Exception as e:
            logger.error("文字提取失败: %s", e)
        return texts

    @staticmethod
    def _extract_slide_text(slide) -> str:
        """提取单页幻灯片的文字内容"""
        parts = []

        for shape in slide.shapes:
            # 文本框
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)

            # 表格
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_texts.append(cell_text)
                    if row_texts:
                        parts.append(" | ".join(row_texts))

        return "\n".join(parts)

    # ============================================================
    # 缩略图
    # ============================================================

    def _create_thumbnail(
        self,
        image_path: str,
        output_dir: Path,
        page_number: int,
    ) -> str:
        """生成缩略图"""
        try:
            thumb_filename = f"thumb_{page_number}.png"
            thumb_path = str(output_dir / thumb_filename)

            with Image.open(image_path) as img:
                ratio = self.thumbnail_width / img.width
                new_height = int(img.height * ratio)
                thumb = img.resize(
                    (self.thumbnail_width, new_height),
                    Image.LANCZOS,
                )
                thumb.save(thumb_path, "PNG")

            return thumb_path
        except Exception as e:
            logger.debug("缩略图生成失败: %s", e)
            return ""

    # ============================================================
    # 文件清理
    # ============================================================

    def delete_ppt_files(self, file_id: str) -> None:
        """删除 PPT 相关的所有文件"""
        file_dir = self.upload_dir / file_id
        if file_dir.exists():
            shutil.rmtree(file_dir, ignore_errors=True)
            logger.info("PPT 文件已清理: %s", file_id)

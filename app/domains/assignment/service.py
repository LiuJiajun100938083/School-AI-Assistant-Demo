#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
作業管理 Service
=================
業務邏輯層，處理作業的創建、提交、批改等完整流程。

用法:
    service = AssignmentService(...)
    # 老師: 創建作業
    assignment = service.create_assignment(teacher_id=1, title="Swift 實作", ...)
    # 學生: 提交作業
    submission = service.submit(assignment_id=1, student=user, files=[...])
    # 老師: AI 批改
    ai_result = service.ai_grade_submission(submission_id=1)
    # 老師: 確認並提交分數
    service.grade_submission(submission_id=1, rubric_scores=[...], feedback="...")
"""

import json
import logging
import os
import re
import uuid
from datetime import datetime
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional

from fastapi import UploadFile

from app.core.exceptions import NotFoundError, ValidationError
from app.domains.assignment.constants import (
    DOCUMENT_EXTENSIONS,
    EXTENSION_TYPE_MAP,
    MAX_FILE_SIZE,
    PREVIEWABLE_ARCHIVE_EXTENSIONS,
    PREVIEWABLE_OFFICE_EXTENSIONS,
    TEXT_READABLE_EXTENSIONS,
)
from app.domains.assignment.exceptions import (
    AssignmentNotFoundError,
    AssignmentNotPublishedError,
    DeadlinePassedError,
    FileTooLargeError,
    InvalidFileTypeError,
    SubmissionNotFoundError,
    TooManyFilesError,
)
from app.domains.assignment.repository import (
    AssignmentAttachmentRepository,
    AssignmentQuestionRepository,
    AssignmentRepository,
    ExamUploadBatchRepository,
    ExamUploadFileRepository,
    RubricItemRepository,
    RubricScoreRepository,
    SubmissionAnswerFileRepository,
    SubmissionAnswerRepository,
    SubmissionFileRepository,
    SubmissionRepository,
)
from app.domains.assignment.schemas import AssignmentType, ScoreSource
from app.domains.user.repository import UserRepository

logger = logging.getLogger(__name__)

# 上傳目錄
UPLOAD_DIR = Path(__file__).resolve().parent.parent.parent.parent / "uploads" / "assignments"


class AssignmentService:
    """作業管理服務"""

    def __init__(
        self,
        assignment_repo: AssignmentRepository,
        submission_repo: SubmissionRepository,
        file_repo: SubmissionFileRepository,
        rubric_repo: RubricItemRepository,
        score_repo: RubricScoreRepository,
        user_repo: UserRepository,
        attachment_repo: Optional["AssignmentAttachmentRepository"] = None,
        conversation_repo=None,
        settings=None,
        question_repo: Optional["AssignmentQuestionRepository"] = None,
        batch_repo: Optional["ExamUploadBatchRepository"] = None,
        upload_file_repo: Optional["ExamUploadFileRepository"] = None,
        answer_repo: Optional["SubmissionAnswerRepository"] = None,
        answer_file_repo: Optional["SubmissionAnswerFileRepository"] = None,
        question_option_repo=None,
    ):
        self._assignment_repo = assignment_repo
        self._submission_repo = submission_repo
        self._file_repo = file_repo
        self._rubric_repo = rubric_repo
        self._score_repo = score_repo
        self._user_repo = user_repo
        self._attachment_repo = attachment_repo
        self._conversation_repo = conversation_repo
        self._settings = settings
        self._question_repo = question_repo
        self._batch_repo = batch_repo
        self._upload_file_repo = upload_file_repo
        self._answer_repo = answer_repo
        self._answer_file_repo = answer_file_repo
        self._question_option_repo = question_option_repo
        self._ask_ai_func: Optional[Callable] = None

        # 確保上傳目錄存在
        UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

    def set_ai_function(self, ask_ai: Callable):
        """注入 AI 函數"""
        self._ask_ai_func = ask_ai

    # ================================================================
    # 老師操作
    # ================================================================

    def _calc_max_score(self, rubric_type: str, rubric_items: Optional[List[Dict]],
                         rubric_config: Optional[Dict] = None) -> Optional[float]:
        """根據評分類型計算滿分"""
        if rubric_type == "competency":
            return None  # 能力等級制無數字分數
        if rubric_type == "holistic":
            if rubric_config and rubric_config.get("levels"):
                return max((lv.get("max", 0) for lv in rubric_config["levels"]), default=100)
            return 100.0
        if rubric_type == "weighted_pct":
            return float(rubric_config.get("total_score", 100)) if rubric_config else 100.0
        if rubric_type == "checklist":
            return float(rubric_config.get("max_score", 100)) if rubric_config else 100.0
        # points / analytic_levels / dse_criterion
        if rubric_items:
            return sum(item.get("max_points", 0) or 0 for item in rubric_items)
        return 100.0

    def create_assignment(
        self,
        teacher_id: int,
        teacher_name: str,
        title: str,
        description: str = "",
        assignment_type: str = AssignmentType.FILE_UPLOAD,
        target_type: str = "all",
        target_value: Optional[str] = None,
        deadline: Optional[str] = None,
        max_files: int = 5,
        allow_late: bool = False,
        rubric_type: str = "points",
        rubric_config: Optional[Dict] = None,
        rubric_items: Optional[List[Dict]] = None,
        questions: Optional[List[Dict]] = None,
        exam_batch_id: Optional[str] = None,
    ) -> Dict[str, Any]:
        """創建作業 (草稿狀態)"""
        if not title or not title.strip():
            raise ValidationError("作業標題不能為空", field="title")

        if assignment_type not in AssignmentType.ALL:
            raise ValidationError(
                f"無效的作業類型: '{assignment_type}'，必須是 {AssignmentType.ALL} 之一",
                field="assignment_type",
            )

        # 解析截止日期
        deadline_dt = None
        if deadline:
            try:
                deadline_dt = datetime.fromisoformat(deadline)
            except ValueError:
                raise ValidationError("截止日期格式無效", field="deadline")

        # Form 類型：滿分 = 題目分數之和
        if assignment_type == AssignmentType.FORM:
            if not questions:
                raise ValidationError("Form 類型作業至少需要 1 道題目", field="questions")
            max_score = sum(q.get("max_points", 0) for q in questions)
        elif assignment_type == AssignmentType.EXAM:
            # Exam 類型：max_score 在保存題目時由 save_assignment_questions 計算
            max_score = 0
        else:
            max_score = self._calc_max_score(rubric_type, rubric_items, rubric_config)

        # 插入作業
        insert_data = {
            "title": title.strip(),
            "description": description.strip() if description else "",
            "created_by": teacher_id,
            "created_by_name": teacher_name,
            "assignment_type": assignment_type,
            "target_type": target_type,
            "target_value": target_value,
            "max_score": max_score,
            "rubric_type": rubric_type,
            "deadline": deadline_dt,
            "status": "draft",
            "allow_late": allow_late,
            "max_files": max_files,
            "created_at": datetime.now(),
            "updated_at": datetime.now(),
        }
        if rubric_config is not None:
            insert_data["rubric_config"] = json.dumps(rubric_config, ensure_ascii=False)
        if exam_batch_id:
            insert_data["exam_batch_id"] = exam_batch_id

        assignment_id = self._assignment_repo.insert_get_id(insert_data)

        # 插入評分標準 (file_upload 類型)
        if assignment_type == AssignmentType.FILE_UPLOAD and rubric_items:
            self._rubric_repo.batch_insert(assignment_id, rubric_items)

        # 插入題目 (試卷識別)
        if questions and self._question_repo:
            self._question_repo.batch_insert(assignment_id, questions)

        logger.info("教師 %s 創建了作業 #%d: %s (類型=%s)", teacher_name, assignment_id, title, rubric_type)
        return self.get_assignment_detail(assignment_id)

    def update_assignment(
        self,
        assignment_id: int,
        teacher_id: int,
        **fields,
    ) -> Dict[str, Any]:
        """更新作業（草稿和已發布狀態均可編輯）"""
        assignment = self._get_assignment_or_raise(assignment_id)

        if assignment["status"] == "closed":
            raise ValidationError("已關閉的作業無法編輯")

        rubric_items = fields.pop("rubric_items", None)
        rubric_type = fields.pop("rubric_type", None)
        rubric_config = fields.pop("rubric_config", None)
        questions = fields.pop("questions", None)

        allowed = {"title", "description", "target_type", "target_value",
                    "deadline", "max_files", "allow_late", "exam_batch_id"}
        update_data = {}
        for key, value in fields.items():
            if key in allowed and value is not None:
                if key == "deadline" and isinstance(value, str):
                    try:
                        update_data[key] = datetime.fromisoformat(value) if value else None
                    except ValueError:
                        raise ValidationError("截止日期格式無效", field="deadline")
                else:
                    update_data[key] = value

        is_form = (assignment.get("assignment_type") or AssignmentType.FILE_UPLOAD) == AssignmentType.FORM
        is_exam = (assignment.get("assignment_type") or AssignmentType.FILE_UPLOAD) == AssignmentType.EXAM

        # Exam 類型：更新題目
        if is_exam and questions is not None and self._question_repo:
            self._question_repo.delete_by_assignment(assignment_id)
            if questions:
                self._question_repo.batch_insert(assignment_id, questions)
                update_data["max_score"] = sum(q.get("points", 0) or 0 for q in questions)

        # Form 類型：題目冻結檢查
        if is_form and questions is not None:
            has_subs = self._submission_repo.count(
                where="assignment_id = %s", params=(assignment_id,)
            ) > 0
            if has_subs:
                raise ValidationError("已有學生提交，無法修改題目結構")
            # 刪除舊題目（CASCADE 清理選項）
            if self._question_repo:
                self._question_repo.delete_by_assignment(assignment_id)
            if questions:
                self._create_form_questions(assignment_id, questions)
            # 重新計算滿分
            update_data["max_score"] = sum(q.get("max_points", 0) for q in questions)

        if rubric_type is not None:
            update_data["rubric_type"] = rubric_type
        if rubric_config is not None:
            update_data["rubric_config"] = json.dumps(rubric_config, ensure_ascii=False)

        # 更新評分標準 (file_upload 類型)
        if not is_form and rubric_items is not None:
            self._rubric_repo.delete_by_assignment(assignment_id)
            if rubric_items:
                self._rubric_repo.batch_insert(assignment_id, rubric_items)
            # 重新計算滿分
            effective_type = rubric_type or assignment.get("rubric_type") or "points"
            effective_config = rubric_config
            if effective_config is None and assignment.get("rubric_config"):
                cfg = assignment["rubric_config"]
                effective_config = json.loads(cfg) if isinstance(cfg, str) else cfg
            update_data["max_score"] = self._calc_max_score(
                effective_type, rubric_items, effective_config
            )

        update_data["updated_at"] = datetime.now()

        if update_data:
            self._assignment_repo.update(
                data=update_data,
                where="id = %s",
                params=(assignment_id,),
            )

        logger.info("教師更新了作業 #%d", assignment_id)
        return self.get_assignment_detail(assignment_id)

    def publish_assignment(self, assignment_id: int, teacher_id: int) -> Dict[str, Any]:
        """發布作業"""
        assignment = self._get_assignment_or_raise(assignment_id)

        if assignment["status"] == "closed":
            raise ValidationError("已關閉的作業不能重新發布")

        self._assignment_repo.update(
            data={
                "status": "published",
                "published_at": datetime.now(),
                "updated_at": datetime.now(),
            },
            where="id = %s AND is_deleted = 0",
            params=(assignment_id,),
        )

        logger.info("作業 #%d 已發布", assignment_id)
        return self.get_assignment_detail(assignment_id)

    def close_assignment(self, assignment_id: int, teacher_id: int) -> Dict[str, Any]:
        """關閉作業 (停止接受提交)"""
        self._get_assignment_or_raise(assignment_id)

        self._assignment_repo.update(
            data={"status": "closed", "updated_at": datetime.now()},
            where="id = %s AND is_deleted = 0",
            params=(assignment_id,),
        )

        logger.info("作業 #%d 已關閉", assignment_id)
        return self.get_assignment_detail(assignment_id)

    def delete_assignment(self, assignment_id: int, teacher_id: int) -> None:
        """軟刪除作業"""
        self._get_assignment_or_raise(assignment_id)
        self._assignment_repo.soft_delete(
            where="id = %s",
            params=(assignment_id,),
        )
        logger.info("作業 #%d 已刪除", assignment_id)

    def list_teacher_assignments(
        self,
        teacher_id: int = 0,
        status: str = "",
        page: int = 1,
        page_size: int = 20,
    ) -> Dict[str, Any]:
        """獲取老師的作業列表 (分頁)"""
        result = self._assignment_repo.find_active(
            status=status,
            created_by=teacher_id,
            page=page,
            page_size=page_size,
        )

        # 為每個作業附加統計和評分標準
        for item in result["items"]:
            stats = self._submission_repo.get_submission_stats(item["id"])
            item["submission_count"] = stats["total_submissions"] if stats else 0
            item["graded_count"] = stats["graded_count"] if stats else 0
            rubric = self._rubric_repo.find_by_assignment(item["id"])
            item["rubric_items"] = rubric
            # OCR 狀態
            self._enrich_ocr_status(item)

        return result

    def _enrich_ocr_status(self, assignment: Dict) -> None:
        """若作業有 exam_batch_id，附加 OCR 狀態資訊"""
        batch_id = assignment.get("exam_batch_id")
        if not batch_id or not self._batch_repo:
            return
        batch = self._batch_repo.find_by_batch_id(batch_id)
        if batch:
            assignment["ocr_status"] = batch.get("status", "unknown")
            # 已完成時附加題目數量
            if batch.get("status") == "completed" and self._question_repo:
                qs = self._question_repo.find_by_assignment(assignment["id"])
                assignment["ocr_question_count"] = len(qs)
            else:
                assignment["ocr_question_count"] = 0
        else:
            assignment["ocr_status"] = "unknown"
            assignment["ocr_question_count"] = 0

    @staticmethod
    def _deserialize_json_fields(data: Dict) -> Dict:
        """反序列化 JSON 字段"""
        for key in ("rubric_config", "level_definitions", "metadata"):
            val = data.get(key)
            if isinstance(val, str):
                try:
                    data[key] = json.loads(val)
                except (json.JSONDecodeError, TypeError):
                    pass
        return data

    def get_assignment_detail(self, assignment_id: int, include_answers: bool = True) -> Dict[str, Any]:
        """獲取作業完整詳情（教師端，含正確答案）"""
        assignment = self._get_assignment_or_raise(assignment_id)
        self._deserialize_json_fields(assignment)

        rubric_items = self._rubric_repo.find_by_assignment(assignment_id)
        for item in rubric_items:
            self._deserialize_json_fields(item)
        assignment["rubric_items"] = rubric_items

        # Form 類型：附帶題目（含正確答案）
        asg_type = assignment.get("assignment_type") or AssignmentType.FILE_UPLOAD
        if asg_type == AssignmentType.FORM:
            assignment["questions"] = self.get_form_questions(
                assignment_id, include_answers=include_answers
            )

        # 附件
        if self._attachment_repo:
            assignment["attachments"] = self._attachment_repo.find_by_assignment(assignment_id)
        else:
            assignment["attachments"] = []

        # 題目
        if self._question_repo:
            questions = self._question_repo.find_by_assignment(assignment_id)
            for q in questions:
                self._deserialize_json_fields(q)
            assignment["questions"] = questions
        else:
            assignment["questions"] = []

        # Form/Exam 類型：從實際題目重新計算滿分
        if asg_type in (AssignmentType.FORM, AssignmentType.EXAM) and assignment.get("questions"):
            assignment["max_score"] = sum(
                float(q.get("max_points") or q.get("points") or 0)
                for q in assignment["questions"]
                if q.get("question_type") != "passage"
            )

        stats = self._submission_repo.get_submission_stats(assignment_id)
        assignment["submission_count"] = stats["total_submissions"] if stats else 0
        assignment["graded_count"] = stats["graded_count"] if stats else 0
        assignment["avg_score"] = float(stats["avg_score"]) if stats and stats["avg_score"] else None

        # OCR 狀態
        self._enrich_ocr_status(assignment)

        return assignment

    def list_submissions(
        self,
        assignment_id: int,
        status: str = "",
    ) -> List[Dict[str, Any]]:
        """獲取某作業的所有提交"""
        self._get_assignment_or_raise(assignment_id)
        submissions = self._submission_repo.find_by_assignment(assignment_id, status)

        # 附加文件和逐項得分
        for sub in submissions:
            sub["files"] = self._file_repo.find_by_submission(sub["id"])
            sub["rubric_scores"] = self._score_repo.find_by_submission(sub["id"])

        return submissions

    def get_submission_detail(self, submission_id: int) -> Dict[str, Any]:
        """獲取單個提交的完整詳情"""
        submission = self._submission_repo.find_by_id(submission_id)
        if not submission:
            raise SubmissionNotFoundError(submission_id)

        submission["files"] = self._file_repo.find_by_submission(submission_id)
        submission["rubric_scores"] = self._score_repo.find_by_submission(submission_id)

        # 附加作業信息 (含評分標準)
        assignment = self._assignment_repo.find_by_id(submission["assignment_id"])
        if assignment:
            self._deserialize_json_fields(assignment)
            submission["assignment"] = assignment
            rubric_items = self._rubric_repo.find_by_assignment(
                submission["assignment_id"]
            )
            for item in rubric_items:
                self._deserialize_json_fields(item)
            submission["rubric_items"] = rubric_items

            # Form / Exam 類型：附帶題目和作答
            atype = assignment.get("assignment_type") or AssignmentType.FILE_UPLOAD
            if atype in (AssignmentType.FORM, AssignmentType.EXAM):
                submission["questions"] = self.get_form_questions(
                    submission["assignment_id"], include_answers=True
                )
                if self._answer_repo:
                    answers = self._answer_repo.find_by_submission(submission_id)
                    if answers and self._answer_file_repo:
                        answer_ids = [a["id"] for a in answers]
                        all_files = self._answer_file_repo.find_by_answers(answer_ids)
                        files_by_answer = {}
                        for f in all_files:
                            files_by_answer.setdefault(f["answer_id"], []).append(f)
                        for a in answers:
                            a["files"] = files_by_answer.get(a["id"], [])
                    submission["answers"] = answers

        return submission

    def _calc_total_score(self, rubric_type: str, rubric_scores: List[Dict],
                           rubric_items: List[Dict], rubric_config: Optional[Dict]) -> Optional[float]:
        """根據評分類型計算學生總分"""
        if rubric_type == "competency":
            return None

        if rubric_type == "holistic":
            # 整體評分: 直接取第一個分數
            if rubric_scores:
                return rubric_scores[0].get("points") or 0
            return 0

        if rubric_type == "checklist":
            # 通過清單: passed / total × max_score
            max_score = float(rubric_config.get("max_score", 100)) if rubric_config else 100.0
            total_items = len(rubric_items) if rubric_items else 1
            passed = sum(1 for s in rubric_scores if (s.get("points") or 0) > 0)
            return round(passed / total_items * max_score, 1)

        if rubric_type == "weighted_pct":
            # 權重百分比: sum(score × weight / 100)
            item_map = {item["id"]: item for item in rubric_items}
            total = 0.0
            for s in rubric_scores:
                item = item_map.get(s["rubric_item_id"], {})
                weight = float(item.get("weight", 0) or 0)
                pts = float(s.get("points", 0) or 0)
                total += pts * weight / 100.0
            return round(total, 1)

        # points / analytic_levels / dse_criterion: sum of points
        return sum(float(s.get("points", 0) or 0) for s in rubric_scores)

    def grade_submission(
        self,
        submission_id: int,
        teacher_id: int,
        rubric_scores: List[Dict],
        feedback: str = "",
    ) -> Dict[str, Any]:
        """批改提交 - 逐項打分 (支持多種評分類型)"""
        submission = self._submission_repo.find_by_id(submission_id)
        if not submission:
            raise SubmissionNotFoundError(submission_id)

        # 獲取作業信息以確定評分類型
        assignment = self._assignment_repo.find_by_id(submission["assignment_id"])
        rubric_type = (assignment or {}).get("rubric_type") or "points"
        rubric_config_raw = (assignment or {}).get("rubric_config")
        rubric_config = json.loads(rubric_config_raw) if isinstance(rubric_config_raw, str) else rubric_config_raw
        rubric_items = self._rubric_repo.find_by_assignment(submission["assignment_id"])

        # 保存逐項得分
        if rubric_scores:
            self._score_repo.batch_upsert(submission_id, rubric_scores)

        # 計算總分
        total_score = self._calc_total_score(rubric_type, rubric_scores, rubric_items, rubric_config)

        # 更新提交狀態
        self._submission_repo.update(
            data={
                "status": "graded",
                "score": total_score,
                "feedback": feedback,
                "graded_by": teacher_id,
                "graded_at": datetime.now(),
                "updated_at": datetime.now(),
            },
            where="id = %s",
            params=(submission_id,),
        )

        logger.info("提交 #%d 已批改，總分: %s (類型=%s)", submission_id, total_score, rubric_type)
        return self.get_submission_detail(submission_id)

    def get_available_targets(self) -> Dict[str, Any]:
        """獲取可選的佈置目標 (班級列表 + 學生列表)"""
        classes_result = self._user_repo.raw_query(
            "SELECT DISTINCT class_name FROM users "
            "WHERE class_name IS NOT NULL AND class_name != '' "
            "ORDER BY class_name"
        )
        classes = [r["class_name"] for r in classes_result]

        students = self._user_repo.find_all(
            where="role = 'student' AND is_active = 1",
            columns="id, username, display_name, class_name",
            order_by="class_name ASC, username ASC",
        )

        return {
            "classes": classes,
            "students": students,
        }

    # ================================================================
    # AI 批改
    # ================================================================

    def _build_ai_prompt(self, rubric_type: str, assignment: Dict,
                          rubric_items: List[Dict], rubric_config: Optional[Dict],
                          submission: Dict, file_contents: str,
                          extra_prompt: str = "") -> str:
        """根據評分類型構建 AI 批改提示"""
        teacher_note = ""
        if extra_prompt and extra_prompt.strip():
            teacher_note = f"""
## 教師批改指示
{extra_prompt.strip()}
（請在批改時遵循以上教師的額外要求）
"""

        base = f"""你是一位專業的作業批改助手。請根據以下評分標準嚴格批改學生的作業。
{teacher_note}
## 作業標題
{assignment['title']}

## 作業描述
{assignment.get('description', '無')}

## 學生提交備註
{submission.get('content', '無')}

## 學生提交文件內容
{file_contents if file_contents else '（無可讀取的文件內容）'}
"""

        if rubric_type == "holistic":
            levels_text = ""
            if rubric_config and rubric_config.get("levels"):
                levels_text = "\n".join(
                    f"- {lv['label']} ({lv.get('min',0)}-{lv.get('max',100)}分): {lv.get('description','')}"
                    for lv in rubric_config["levels"]
                )
            return base + f"""
## 整體評分等級
{levels_text}

請以嚴格 JSON 格式返回:
{{"selected_level": "等級標籤", "points": 分數, "reason": "評分理由", "overall_feedback": "總體評語"}}
"""

        if rubric_type == "checklist":
            items_text = "\n".join(f"- id={item['id']}: {item['title']}" for item in rubric_items)
            return base + f"""
## 檢查清單 (每項判斷通過/不通過)
{items_text}

請以嚴格 JSON 格式返回:
{{"items": [{{"rubric_item_id": ID, "passed": true或false, "reason": "理由"}}, ...], "overall_feedback": "總體評語"}}
"""

        if rubric_type == "competency":
            level_labels = (rubric_config or {}).get("level_labels", ["Not Yet", "Approaching", "Meeting", "Exceeding"])
            labels_text = " / ".join(level_labels)
            items_text = "\n".join(f"- id={item['id']}: {item['title']}" for item in rubric_items)
            return base + f"""
## 能力等級評估 (等級: {labels_text})
{items_text}

請以嚴格 JSON 格式返回 (無數字分數):
{{"items": [{{"rubric_item_id": ID, "selected_level": "等級", "reason": "理由"}}, ...], "overall_feedback": "總體評語"}}
"""

        if rubric_type == "analytic_levels":
            items_text = ""
            for item in rubric_items:
                ld = item.get("level_definitions") or []
                if isinstance(ld, str):
                    try: ld = json.loads(ld)
                    except: ld = []
                levels = ", ".join(f"{l['level']}={l.get('points',0)}分" for l in ld)
                items_text += f"- id={item['id']}: {item['title']} (滿分{item['max_points']}) 等級: [{levels}]\n"
            return base + f"""
## 分級評分標準
{items_text}

請為每項選擇一個等級，以嚴格 JSON 格式返回:
{{"items": [{{"rubric_item_id": ID, "selected_level": "等級名", "points": 分數, "reason": "理由"}}, ...], "overall_feedback": "總體評語"}}
"""

        if rubric_type == "weighted_pct":
            total_score = (rubric_config or {}).get("total_score", 100)
            items_text = "\n".join(
                f"- id={item['id']}: {item['title']} (權重{item.get('weight',0)}%)"
                for item in rubric_items
            )
            return base + f"""
## 權重百分比評分 (滿分 {total_score})
{items_text}

請為每項打分 (0-{total_score})，以嚴格 JSON 格式返回:
{{"items": [{{"rubric_item_id": ID, "points": 分數, "reason": "理由"}}, ...], "overall_feedback": "總體評語"}}
"""

        if rubric_type == "dse_criterion":
            items_text = ""
            for item in rubric_items:
                ld = item.get("level_definitions") or []
                if isinstance(ld, str):
                    try: ld = json.loads(ld)
                    except: ld = []
                desc = "; ".join(f"Level {l['level']}: {l.get('description','')}" for l in ld if l.get("description"))
                items_text += f"- id={item['id']}: {item['title']} (max {item['max_points']}) [{desc}]\n"
            return base + f"""
## DSE 標準量規
{items_text}

請為每項選擇等級 (0-max)，以嚴格 JSON 格式返回:
{{"items": [{{"rubric_item_id": ID, "points": 等級數, "reason": "理由"}}, ...], "overall_feedback": "總體評語"}}
"""

        # Default: points
        rubric_text = "\n".join(
            f"{i+1}. {item['title']} (滿分 {item['max_points']} 分)"
            for i, item in enumerate(rubric_items)
        )
        id_hint = ", ".join(
            f'id={item["id"]}:"{item["title"]}"(滿分{item["max_points"]})'
            for item in rubric_items
        )
        return base + f"""
## 評分標準
{rubric_text}

請以嚴格的 JSON 格式返回批改結果，不要包含任何其他文字：
{{"items": [{{"rubric_item_id": {rubric_items[0]['id']}, "points": 分數, "reason": "評分理由"}}, ...], "overall_feedback": "總體評語"}}

注意：
- 每個 rubric_item_id 必須對應上述評分標準的 ID
- points 不能超過該項的滿分
- reason 請簡短說明得分或失分原因
- overall_feedback 請給出總體評價和改進建議

評分標準 ID 對照: {id_hint}"""

    def ai_grade_submission(self, submission_id: int, extra_prompt: str = "") -> Dict[str, Any]:
        """AI 自動批改一份提交"""
        if not self._ask_ai_func:
            raise ValidationError("AI 批改功能未初始化")

        submission = self._submission_repo.find_by_id(submission_id)
        if not submission:
            raise SubmissionNotFoundError(submission_id)

        assignment = self._assignment_repo.find_by_id(submission["assignment_id"])
        if not assignment:
            raise AssignmentNotFoundError(submission["assignment_id"])
        self._deserialize_json_fields(assignment)

        rubric_type = assignment.get("rubric_type") or "points"
        rubric_config = assignment.get("rubric_config")

        # 讀取評分標準
        rubric_items = self._rubric_repo.find_by_assignment(assignment["id"])
        for item in rubric_items:
            self._deserialize_json_fields(item)

        if not rubric_items and rubric_type != "holistic":
            raise ValidationError("此作業沒有設定評分標準，無法使用 AI 批改")

        # 讀取提交文件
        files = self._file_repo.find_by_submission(submission_id)
        file_contents = self._extract_file_contents(files)

        prompt = self._build_ai_prompt(
            rubric_type, assignment, rubric_items, rubric_config,
            submission, file_contents, extra_prompt=extra_prompt
        )

        try:
            model = None
            if self._settings and hasattr(self._settings, 'llm_local_model'):
                model = self._settings.llm_local_model

            answer, _ = self._ask_ai_func(
                question=prompt,
                subject_code="general",
                use_api=False,
                conversation_history=[],
                model=model,
            )

            result = self._parse_ai_response(answer, rubric_items, rubric_type)
            return result

        except Exception as e:
            logger.error("AI 批改失敗: %s", e)
            return {
                "items": [],
                "overall_feedback": f"AI 批改失敗: {str(e)}",
                "error": True,
            }

    def _parse_ai_response(
        self, answer: str, rubric_items: List[Dict], rubric_type: str = "points"
    ) -> Dict[str, Any]:
        """解析 AI 回應中的 JSON（含截斷修復）"""
        # 1) 直接解析
        try:
            result = json.loads(answer)
            return self._validate_ai_result(result, rubric_items, rubric_type)
        except json.JSONDecodeError:
            pass

        # 2) 提取 JSON 區塊
        json_match = re.search(r'\{[\s\S]*\}', answer)
        if json_match:
            try:
                result = json.loads(json_match.group())
                return self._validate_ai_result(result, rubric_items, rubric_type)
            except json.JSONDecodeError:
                pass

        # 3) 嘗試修復截斷的 JSON — 逐個提取完整的 item 對象
        repaired = self._repair_truncated_json(answer, rubric_type)
        if repaired:
            return self._validate_ai_result(repaired, rubric_items, rubric_type)

        return {
            "items": [],
            "overall_feedback": answer[:500],
            "error": True,
        }

    def _repair_truncated_json(self, answer: str, rubric_type: str) -> Optional[Dict]:
        """修復被截斷的 AI JSON 回應，提取已完成的 item"""
        try:
            # 提取所有完整的 item 對象 {...}
            item_pattern = r'\{\s*"rubric_item_id"\s*:\s*\d+[^{}]*\}'
            items = re.findall(item_pattern, answer)
            if not items:
                return None

            parsed_items = []
            for item_str in items:
                try:
                    parsed_items.append(json.loads(item_str))
                except json.JSONDecodeError:
                    continue

            if not parsed_items:
                return None

            # 提取 overall_feedback（如果有）
            fb_match = re.search(r'"overall_feedback"\s*:\s*"([^"]*)"', answer)
            feedback = fb_match.group(1) if fb_match else "（AI 回應被截斷，已自動修復部分評分）"

            logger.warning("AI 回應被截斷，成功修復 %d/%d 項評分", len(parsed_items), len(items))
            return {
                "items": parsed_items,
                "overall_feedback": feedback,
            }
        except Exception:
            return None

    def _validate_ai_result(
        self, result: Dict, rubric_items: List[Dict], rubric_type: str = "points"
    ) -> Dict[str, Any]:
        """驗證和修正 AI 結果"""
        # 整體評分特殊處理
        if rubric_type == "holistic":
            return {
                "selected_level": result.get("selected_level", ""),
                "points": float(result.get("points", 0)),
                "reason": result.get("reason", ""),
                "overall_feedback": result.get("overall_feedback", ""),
                "items": [],
                "error": False,
            }

        valid_ids = {item["id"]: item for item in rubric_items}
        validated_items = []

        for item in result.get("items", []):
            rid = item.get("rubric_item_id")
            if rid not in valid_ids:
                continue

            rubric_item = valid_ids[rid]
            entry = {
                "rubric_item_id": rid,
                "reason": item.get("reason", ""),
            }

            if rubric_type == "checklist":
                passed = item.get("passed", False)
                entry["points"] = 1 if passed else 0
                entry["passed"] = passed
            elif rubric_type == "competency":
                entry["selected_level"] = item.get("selected_level", "")
                entry["points"] = None
            else:
                points = float(item.get("points", 0) or 0)
                max_p = float(rubric_item.get("max_points", 0) or 1000)
                points = max(0, min(points, max_p))
                entry["points"] = points
                if item.get("selected_level"):
                    entry["selected_level"] = item["selected_level"]

            validated_items.append(entry)

        return {
            "items": validated_items,
            "overall_feedback": result.get("overall_feedback", ""),
            "error": False,
        }

    def _extract_file_contents(self, files: List[Dict]) -> str:
        """提取文件的文本內容 (用於 AI 批改)"""
        contents = []

        for f in files:
            file_path = UPLOAD_DIR / f.get("stored_name", "")
            original_name = f.get("original_name", "")
            ext = Path(original_name).suffix.lower()

            if ext in TEXT_READABLE_EXTENSIONS:
                # 直接讀取代碼/文本文件
                try:
                    if file_path.exists():
                        encodings = ['utf-8', 'gbk', 'gb2312', 'big5', 'utf-16']
                        text = None
                        for encoding in encodings:
                            try:
                                with open(file_path, 'r', encoding=encoding) as fh:
                                    text = fh.read()
                                break
                            except (UnicodeDecodeError, UnicodeError):
                                continue

                        if text:
                            # 清除 Xcode/Swift Playgrounds 內部佔位標記
                            text = re.sub(r'/\*@[A-Z_]+@\*/', '', text)
                            # 限制每個文件的文本長度
                            if len(text) > 10000:
                                text = text[:10000] + "\n... (文件過長，已截斷)"
                            contents.append(
                                f"### 文件: {original_name}\n```\n{text}\n```"
                            )
                except Exception as e:
                    logger.warning("讀取文件 %s 失敗: %s", original_name, e)

            elif ext in DOCUMENT_EXTENSIONS:
                # 嘗試用 FileProcessor 提取
                try:
                    from llm.rag.file_processor import FileProcessor
                    processor = FileProcessor()
                    success, text, _ = processor.process_file(
                        str(file_path), original_name
                    )
                    if success and text:
                        if len(text) > 10000:
                            text = text[:10000] + "\n... (文件過長，已截斷)"
                        contents.append(
                            f"### 文件: {original_name}\n{text}"
                        )
                except Exception as e:
                    logger.warning("處理文件 %s 失敗: %s", original_name, e)
                    contents.append(f"### 文件: {original_name}\n（無法提取內容）")
            elif ext in (".swiftpm", ".zip"):
                # 解壓壓縮包，提取源碼（支持 .swiftpm 和 .zip）
                try:
                    import zipfile
                    if file_path.exists() and zipfile.is_zipfile(file_path):
                        with zipfile.ZipFile(file_path, "r") as zf:
                            code_exts = {".swift", ".py", ".js", ".ts", ".java", ".c", ".cpp", ".h", ".html", ".css", ".json"}
                            swift_entries = sorted([
                                n for n in zf.namelist()
                                if any(n.endswith(ce) for ce in code_exts)
                                and not n.startswith("__MACOSX/")
                                and "/.build/" not in n
                                and not n.split("/")[-1].startswith(".")
                            ])
                            remaining = 10000
                            for entry_name in swift_entries:
                                if remaining <= 0:
                                    break
                                try:
                                    raw = zf.read(entry_name)
                                    text = raw.decode("utf-8", errors="replace")
                                except Exception:
                                    continue
                                text = re.sub(r'/\*@[A-Z_]+@\*/', '', text)
                                if len(text) > remaining:
                                    text = text[:remaining] + "\n... (已截斷)"
                                display_name = "/".join(entry_name.rsplit("/", 2)[-2:]) if "/" in entry_name else entry_name
                                contents.append(
                                    f"### 文件: {original_name} → {display_name}\n```swift\n{text}\n```"
                                )
                                remaining -= len(text)
                    else:
                        contents.append(f"### 文件: {original_name}\n（無法解壓文件）")
                except Exception as e:
                    logger.warning("解壓 %s 失敗: %s", original_name, e)
                    contents.append(f"### 文件: {original_name}\n（解壓失敗）")
            else:
                contents.append(f"### 文件: {original_name}\n（{f.get('file_type', '未知')} 類型，無法提取文本）")

        return "\n\n".join(contents)

    # ================================================================
    # 學生操作
    # ================================================================

    def list_student_assignments(
        self,
        username: str,
        class_name: str = "",
        status_filter: str = "",
    ) -> List[Dict[str, Any]]:
        """獲取學生可見的作業列表"""
        assignments = self._assignment_repo.find_published_for_student(
            username, class_name
        )

        # 查詢學生的用戶 ID
        user = self._user_repo.find_one(
            where="username = %s",
            params=(username,),
            columns="id",
        )
        student_id = user["id"] if user else 0

        result = []
        for a in assignments:
            # 查詢提交狀態
            submission = self._submission_repo.find_by_assignment_student(
                a["id"], student_id
            )
            submission_status = submission["status"] if submission else "not_submitted"
            score = submission["score"] if submission and submission.get("score") is not None else None

            # 篩選
            if status_filter == "not_submitted" and submission:
                continue
            if status_filter == "submitted" and (not submission or submission["status"] != "submitted"):
                continue
            if status_filter == "graded" and (not submission or submission["status"] != "graded"):
                continue

            a["submission_status"] = submission_status
            a["my_score"] = score
            a["rubric_items"] = self._rubric_repo.find_by_assignment(a["id"])
            result.append(a)

        return result

    def get_student_assignment_detail(
        self,
        assignment_id: int,
        username: str,
    ) -> Dict[str, Any]:
        """獲取學生的作業詳情 (含自己的提交)"""
        assignment = self._get_assignment_or_raise(assignment_id)

        if assignment["status"] != "published" and assignment["status"] != "closed":
            raise AssignmentNotPublishedError()

        assignment["rubric_items"] = self._rubric_repo.find_by_assignment(assignment_id)

        # Form 類型：附帶題目（不含正確答案和參考答案）
        asg_type = assignment.get("assignment_type") or AssignmentType.FILE_UPLOAD
        is_form = asg_type == AssignmentType.FORM
        is_exam = asg_type == AssignmentType.EXAM
        if is_form:
            assignment["questions"] = self.get_form_questions(
                assignment_id, include_answers=False
            )
        elif is_exam:
            assignment["questions"] = self.get_assignment_questions_for_student(
                assignment_id
            )

        # 從實際題目重新計算滿分（防止 DB 中 max_score 不準確）
        if (is_form or is_exam) and assignment.get("questions"):
            assignment["max_score"] = sum(
                float(q.get("max_points") or q.get("points") or 0)
                for q in assignment["questions"]
                if q.get("question_type") != "passage"
            )

        # 附件
        if self._attachment_repo:
            assignment["attachments"] = self._attachment_repo.find_by_assignment(assignment_id)
        else:
            assignment["attachments"] = []

        # 查詢學生的提交
        user = self._user_repo.find_one(
            where="username = %s",
            params=(username,),
            columns="id",
        )
        student_id = user["id"] if user else 0

        submission = self._submission_repo.find_by_assignment_student(
            assignment_id, student_id
        )
        if submission:
            submission["files"] = self._file_repo.find_by_submission(submission["id"])
            submission["rubric_scores"] = self._score_repo.find_by_submission(submission["id"])
            # Exam 類型：附帶作答詳情
            if is_exam and self._answer_repo:
                answers = self._answer_repo.find_by_submission(submission["id"])
                submission["answers"] = answers or []
            # Form 類型：附帶作答詳情
            if is_form and self._answer_repo:
                answers = self._answer_repo.find_by_submission(submission["id"])
                # 附帶作答文件
                if answers and self._answer_file_repo:
                    answer_ids = [a["id"] for a in answers]
                    all_files = self._answer_file_repo.find_by_answers(answer_ids)
                    files_by_answer = {}
                    for f in all_files:
                        files_by_answer.setdefault(f["answer_id"], []).append(f)
                    for a in answers:
                        a["files"] = files_by_answer.get(a["id"], [])
                submission["answers"] = answers
                # 已提交後，學生可以看到 MC 的正確答案
                if assignment.get("questions"):
                    for q in assignment["questions"]:
                        if q["question_type"] == "mc":
                            # 從完整題目中取回正確答案
                            full_questions = self._question_repo.find_by_assignment(assignment_id) if self._question_repo else []
                            fq_map = {fq["id"]: fq for fq in full_questions}
                            if q["id"] in fq_map:
                                q["correct_answer"] = fq_map[q["id"]].get("correct_answer")

        assignment["my_submission"] = submission
        return assignment

    async def submit_assignment(
        self,
        assignment_id: int,
        student: Dict,
        content: str = "",
        files: Optional[List[UploadFile]] = None,
        is_teacher_proxy: bool = False,
    ) -> Dict[str, Any]:
        """學生提交作業（is_teacher_proxy=True 時跳過狀態與截止日期檢查）"""
        assignment = self._get_assignment_or_raise(assignment_id)

        if not is_teacher_proxy:
            if assignment["status"] != "published":
                raise AssignmentNotPublishedError()

            # 檢查截止日期
            if assignment.get("deadline"):
                deadline = assignment["deadline"]
                if isinstance(deadline, str):
                    deadline = datetime.fromisoformat(deadline)
                if datetime.now() > deadline and not assignment.get("allow_late"):
                    raise DeadlinePassedError()

        # 檢查文件數量
        max_files = assignment.get("max_files", 5)
        if files and len(files) > max_files:
            raise TooManyFilesError(max_files)

        student_id = student["id"]
        is_late = False
        if assignment.get("deadline"):
            deadline = assignment["deadline"]
            if isinstance(deadline, str):
                deadline = datetime.fromisoformat(deadline)
            is_late = datetime.now() > deadline

        # 檢查是否已提交 (如已提交則更新)
        existing = self._submission_repo.find_by_assignment_student(
            assignment_id, student_id
        )

        if existing:
            # 更新現有提交
            submission_id = existing["id"]
            self._submission_repo.update(
                data={
                    "content": content,
                    "status": "submitted",
                    "is_late": is_late,
                    "score": None,
                    "feedback": None,
                    "graded_by": None,
                    "graded_at": None,
                    "updated_at": datetime.now(),
                },
                where="id = %s",
                params=(submission_id,),
            )
            # 刪除舊文件記錄 (物理文件暫保留)
            self._file_repo.delete_by_submission(submission_id)
            self._score_repo.delete_by_submission(submission_id)
        else:
            # 創建新提交
            submission_id = self._submission_repo.insert_get_id({
                "assignment_id": assignment_id,
                "student_id": student_id,
                "student_name": student.get("display_name", ""),
                "username": student.get("username", ""),
                "class_name": student.get("class_name", ""),
                "content": content,
                "status": "submitted",
                "is_late": is_late,
                "submitted_at": datetime.now(),
                "updated_at": datetime.now(),
            })

        # 保存文件
        if files:
            for f in files:
                await self._save_upload_file(submission_id, f)

        logger.info(
            "學生 %s 提交了作業 #%d",
            student.get("username"), assignment_id,
        )

        return self.get_submission_detail(submission_id)

    async def _save_upload_file(
        self, submission_id: int, file: UploadFile
    ) -> Dict[str, Any]:
        """保存上傳的文件"""
        original_name = file.filename or "unnamed"
        ext = Path(original_name).suffix.lower()

        # 檢查文件類型
        file_type = EXTENSION_TYPE_MAP.get(ext)
        if not file_type:
            raise InvalidFileTypeError(ext)

        # 讀取文件內容
        content = await file.read()
        file_size = len(content)

        # 檢查大小
        if file_size > MAX_FILE_SIZE:
            raise FileTooLargeError(MAX_FILE_SIZE // 1024 // 1024)

        # UUID 命名存儲
        stored_name = f"{uuid.uuid4().hex}{ext}"
        file_path = UPLOAD_DIR / stored_name

        with open(file_path, "wb") as fh:
            fh.write(content)

        # 保存數據庫記錄
        file_id = self._file_repo.insert_get_id({
            "submission_id": submission_id,
            "original_name": original_name,
            "stored_name": stored_name,
            "file_path": f"uploads/assignments/{stored_name}",
            "file_size": file_size,
            "file_type": file_type,
            "mime_type": file.content_type or "",
        })

        return {
            "id": file_id,
            "original_name": original_name,
            "stored_name": stored_name,
            "file_size": file_size,
            "file_type": file_type,
        }

    # ================================================================
    # 作業附件
    # ================================================================

    async def upload_attachment(
        self, assignment_id: int, teacher_id: int, file: UploadFile
    ) -> Dict[str, Any]:
        """上傳作業附件（教師用）"""
        assignment = self._get_assignment_or_raise(assignment_id)
        if assignment.get("created_by") != teacher_id:
            from app.core.exceptions import AuthorizationError
            raise AuthorizationError("只有作業創建者可以上傳附件")

        if not self._attachment_repo:
            raise ValidationError("附件功能未初始化")

        original_name = file.filename or "unnamed"
        ext = Path(original_name).suffix.lower()

        file_type = EXTENSION_TYPE_MAP.get(ext)
        if not file_type:
            raise InvalidFileTypeError(ext)

        content = await file.read()
        file_size = len(content)
        if file_size > MAX_FILE_SIZE:
            raise FileTooLargeError(MAX_FILE_SIZE // 1024 // 1024)

        stored_name = f"{uuid.uuid4().hex}{ext}"
        file_path = UPLOAD_DIR / stored_name

        with open(file_path, "wb") as fh:
            fh.write(content)

        file_id = self._attachment_repo.insert_get_id({
            "assignment_id": assignment_id,
            "original_name": original_name,
            "stored_name": stored_name,
            "file_path": f"uploads/assignments/{stored_name}",
            "file_size": file_size,
            "file_type": file_type,
            "mime_type": file.content_type or "",
        })

        logger.info("作業 #%d 上傳附件: %s (%d bytes)", assignment_id, original_name, file_size)
        return {
            "id": file_id,
            "original_name": original_name,
            "stored_name": stored_name,
            "file_path": f"uploads/assignments/{stored_name}",
            "file_size": file_size,
            "file_type": file_type,
        }

    def delete_attachment(self, attachment_id: int, teacher_id: int) -> bool:
        """刪除作業附件（軟刪除）"""
        if not self._attachment_repo:
            raise ValidationError("附件功能未初始化")

        attachment = self._attachment_repo.find_by_id(attachment_id)
        if not attachment:
            raise NotFoundError("附件不存在")

        # 驗證歸屬
        assignment = self._assignment_repo.find_by_id(attachment["assignment_id"])
        if not assignment or assignment.get("created_by") != teacher_id:
            from app.core.exceptions import AuthorizationError
            raise AuthorizationError("只有作業創建者可以刪除附件")

        self._attachment_repo.soft_delete_attachment(attachment_id)
        logger.info("刪除附件 #%d (作業 #%d)", attachment_id, attachment["assignment_id"])
        return True

    # ================================================================
    # 文件預覽（Office → HTML）
    # ================================================================

    _preview_cache: Dict[str, Dict[str, Any]] = {}  # class-level 簡易緩存

    def preview_file(self, file_id: int, user_id: int) -> Dict[str, Any]:
        """
        將 docx/xlsx/pptx 轉為 HTML 預覽。
        帶緩存：file_id + file_size 不變即命中。
        """
        import html as html_mod

        file_rec = self._file_repo.find_by_id(file_id)
        if not file_rec:
            raise NotFoundError("文件不存在")

        # ---- 權限校驗 ----
        submission = self._submission_repo.find_by_id(file_rec["submission_id"])
        if not submission:
            raise NotFoundError("提交不存在")
        assignment = self._assignment_repo.find_by_id(submission["assignment_id"])
        if not assignment:
            raise NotFoundError("作業不存在")
        # 老師（作業創建者）或學生本人可預覽
        is_owner = submission.get("student_id") == user_id
        is_teacher = assignment.get("created_by") == user_id
        if not is_owner and not is_teacher:
            from app.core.exceptions import AuthorizationError
            raise AuthorizationError("無權預覽此文件")

        # ---- 類型白名單 ----
        original_name = file_rec.get("original_name", "")
        ext = Path(original_name).suffix.lower()
        if ext not in PREVIEWABLE_OFFICE_EXTENSIONS and ext not in PREVIEWABLE_ARCHIVE_EXTENSIONS:
            raise ValidationError(f"不支持預覽此文件類型: {ext}")

        # ---- 緩存 ----
        cache_key = f"preview:{file_id}:{file_rec.get('file_size', 0)}"
        if cache_key in self._preview_cache:
            return self._preview_cache[cache_key]

        # ---- 讀取文件 ----
        stored_name = file_rec.get("stored_name", "")
        file_path = UPLOAD_DIR / stored_name
        if not file_path.exists():
            raise NotFoundError("文件不存在於磁盤")

        # ---- 轉換 ----
        try:
            if ext == ".docx":
                result = self._preview_docx(file_path, file_id, html_mod)
            elif ext == ".xlsx":
                result = self._preview_xlsx(file_path, html_mod)
            elif ext == ".pptx":
                result = self._preview_pptx(file_path, file_id, html_mod)
            elif ext in PREVIEWABLE_ARCHIVE_EXTENSIONS:
                result = self._preview_swiftpm(file_path, html_mod)
            else:
                raise ValidationError(f"不支持預覽: {ext}")
        except (ValidationError, NotFoundError):
            raise
        except Exception as e:
            logger.error("預覽文件 #%d 失敗: %s", file_id, e)
            raise ValidationError(f"文件預覽失敗: {str(e)[:200]}")

        result["success"] = True
        result["file_type"] = ext.lstrip(".")

        # 寫入緩存（最多保留 200 條）
        if len(self._preview_cache) > 200:
            # 簡單 FIFO：刪除最早的 50 條
            keys = list(self._preview_cache.keys())[:50]
            for k in keys:
                del self._preview_cache[k]
        self._preview_cache[cache_key] = result
        return result

    # ---- docx → HTML ----
    def _preview_docx(self, file_path: Path, file_id: int, html_mod) -> Dict[str, Any]:
        from docx import Document
        from docx.oxml.ns import qn

        doc = Document(str(file_path))
        parts = []
        preview_dir = UPLOAD_DIR / "previews"
        preview_dir.mkdir(exist_ok=True)
        img_count = 0

        for element in doc.element.body:
            tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

            if tag == "p":
                # 段落
                para_html = self._docx_para_to_html(element, doc, file_id, preview_dir, html_mod)
                if para_html:
                    parts.append(para_html)
                    # 統計圖片
                    img_count += para_html.count("<img ")

            elif tag == "tbl":
                # 表格
                tbl_html = self._docx_table_to_html(element, html_mod)
                parts.append(tbl_html)

        return {
            "html": f'<div class="doc-preview">{"".join(parts)}</div>',
            "truncated": False,
            "meta": {"paragraphs": len(parts), "images": img_count},
        }

    def _docx_para_to_html(self, para_el, doc, file_id, preview_dir, html_mod) -> str:
        from docx.oxml.ns import qn

        runs_html = []
        for run_el in para_el.iter(qn("w:r")):
            # 檢查是否有圖片
            drawings = run_el.findall(f".//{qn('a:blip')}")
            for blip in drawings:
                r_id = blip.get(qn("r:embed"))
                if r_id:
                    img_url = self._extract_docx_image(doc, r_id, file_id, preview_dir)
                    if img_url:
                        runs_html.append(f'<img src="{html_mod.escape(img_url)}" style="max-width:100%;height:auto;">')

            # 文字
            for t_el in run_el.findall(qn("w:t")):
                text = t_el.text or ""
                if not text:
                    continue
                # 檢查粗體/斜體
                rPr = run_el.find(qn("w:rPr"))
                bold = rPr is not None and rPr.find(qn("w:b")) is not None
                italic = rPr is not None and rPr.find(qn("w:i")) is not None
                escaped = html_mod.escape(text)
                if bold:
                    escaped = f"<strong>{escaped}</strong>"
                if italic:
                    escaped = f"<em>{escaped}</em>"
                runs_html.append(escaped)

        if not runs_html:
            return ""

        # 檢查對齊
        pPr = para_el.find(qn("w:pPr"))
        align = ""
        if pPr is not None:
            jc = pPr.find(qn("w:jc"))
            if jc is not None:
                val = jc.get(qn("w:val"), "")
                if val in ("center", "right", "both"):
                    align_map = {"center": "center", "right": "right", "both": "justify"}
                    align = f' style="text-align:{align_map.get(val, "left")}"'

        # 檢查是否標題
        style_el = pPr.find(qn("w:pStyle")) if pPr is not None else None
        style_val = style_el.get(qn("w:val"), "") if style_el is not None else ""
        heading_map = {"Heading1": "h1", "Heading2": "h2", "Heading3": "h3",
                       "Heading4": "h4", "1": "h1", "2": "h2", "3": "h3", "4": "h4"}
        tag = heading_map.get(style_val, "p")

        return f'<{tag}{align}>{"".join(runs_html)}</{tag}>'

    def _extract_docx_image(self, doc, r_id, file_id, preview_dir) -> str:
        """從 docx 提取圖片到 previews/ 目錄，返回 URL"""
        try:
            part = doc.part.related_parts.get(r_id)
            if not part:
                return ""
            img_data = part.blob
            # 判斷格式
            ct = part.content_type or ""
            ext_map = {"image/png": ".png", "image/jpeg": ".jpg", "image/gif": ".gif",
                       "image/bmp": ".bmp", "image/svg+xml": ".svg", "image/webp": ".webp"}
            img_ext = ext_map.get(ct, ".png")
            img_name = f"docx_{file_id}_{r_id}{img_ext}"
            img_path = preview_dir / img_name
            if not img_path.exists():
                with open(img_path, "wb") as f:
                    f.write(img_data)
            return f"/uploads/assignments/previews/{img_name}"
        except Exception as e:
            logger.warning("提取 docx 圖片失敗: %s", e)
            return ""

    def _docx_table_to_html(self, tbl_el, html_mod) -> str:
        from docx.oxml.ns import qn
        rows_html = []
        for tr in tbl_el.findall(qn("w:tr")):
            cells = []
            for tc in tr.findall(qn("w:tc")):
                cell_text = ""
                for p in tc.findall(qn("w:p")):
                    for t in p.iter(qn("w:t")):
                        cell_text += t.text or ""
                cells.append(f"<td>{html_mod.escape(cell_text)}</td>")
            rows_html.append(f"<tr>{''.join(cells)}</tr>")
        return f'<table class="doc-table">{"".join(rows_html)}</table>'

    # ---- xlsx → HTML ----
    def _preview_xlsx(self, file_path: Path, html_mod) -> Dict[str, Any]:
        from openpyxl import load_workbook

        wb = load_workbook(str(file_path), read_only=True, data_only=True)
        sheets_html = []
        truncated = False
        total_rows = 0
        total_cols = 0
        max_rows = 200
        max_cols = 50

        for sheet_name in wb.sheetnames[:10]:  # 最多 10 個 sheet
            ws = wb[sheet_name]
            rows_html = []
            row_count = 0
            col_count = 0

            for row in ws.iter_rows(max_row=max_rows + 1, max_col=max_cols + 1, values_only=False):
                row_count += 1
                if row_count > max_rows:
                    truncated = True
                    break
                cells = []
                col_idx = 0
                for cell in row:
                    col_idx += 1
                    if col_idx > max_cols:
                        truncated = True
                        break
                    val = cell.value
                    text = html_mod.escape(str(val)) if val is not None else ""
                    tag = "th" if row_count == 1 else "td"
                    cells.append(f"<{tag}>{text}</{tag}>")
                col_count = max(col_count, col_idx)
                rows_html.append(f"<tr>{''.join(cells)}</tr>")

            total_rows = max(total_rows, row_count)
            total_cols = max(total_cols, col_count)
            sheet_title = html_mod.escape(sheet_name)
            sheets_html.append(
                f'<div class="sheet-section">'
                f'<h4 class="sheet-title">{sheet_title}</h4>'
                f'<div class="sheet-table-wrapper"><table class="doc-table">{"".join(rows_html)}</table></div>'
                f'</div>'
            )

        wb.close()

        return {
            "html": f'<div class="doc-preview xlsx-preview">{"".join(sheets_html)}</div>',
            "truncated": truncated,
            "meta": {
                "sheets": len(sheets_html),
                "rendered_rows": min(total_rows, max_rows),
                "rendered_cols": min(total_cols, max_cols),
            },
        }

    # ---- pptx → HTML ----
    def _preview_pptx(self, file_path: Path, file_id: int, html_mod) -> Dict[str, Any]:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(str(file_path))
        preview_dir = UPLOAD_DIR / "previews"
        preview_dir.mkdir(exist_ok=True)
        slides_html = []
        max_slides = 50
        img_idx = 0

        for i, slide in enumerate(prs.slides):
            if i >= max_slides:
                break
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            # 標題形狀用 h3，其他用 p
                            tag = "h3" if shape.shape_id == slide.shapes.title.shape_id else "p" if slide.shapes.title and hasattr(slide.shapes, 'title') else "p"
                            try:
                                tag = "h3" if shape == slide.shapes.title else "p"
                            except Exception:
                                tag = "p"
                            parts.append(f"<{tag}>{html_mod.escape(text)}</{tag}>")

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        img_ext = {"image/png": ".png", "image/jpeg": ".jpg",
                                   "image/gif": ".gif"}.get(image.content_type, ".png")
                        img_name = f"pptx_{file_id}_{i}_{img_idx}{img_ext}"
                        img_path = preview_dir / img_name
                        if not img_path.exists():
                            with open(img_path, "wb") as f:
                                f.write(image.blob)
                        img_url = f"/uploads/assignments/previews/{img_name}"
                        parts.append(f'<img src="{html_mod.escape(img_url)}" style="max-width:100%;height:auto;">')
                        img_idx += 1
                    except Exception as e:
                        logger.warning("提取 pptx 圖片失敗: %s", e)

                if shape.has_table:
                    tbl = shape.table
                    rows_html = []
                    for row in tbl.rows:
                        cells = [f"<td>{html_mod.escape(cell.text)}</td>" for cell in row.cells]
                        rows_html.append(f"<tr>{''.join(cells)}</tr>")
                    parts.append(f'<table class="doc-table">{"".join(rows_html)}</table>')

            slide_num = i + 1
            slide_content = "".join(parts) if parts else f'<p style="color:var(--text-tertiary);">（第 {slide_num} 頁無文字內容）</p>'
            slides_html.append(
                f'<div class="slide-card">'
                f'<div class="slide-num">第 {slide_num} 頁</div>'
                f'{slide_content}'
                f'</div>'
            )

        return {
            "html": f'<div class="doc-preview pptx-preview">{"".join(slides_html)}</div>',
            "truncated": len(prs.slides) > max_slides,
            "meta": {"pages": len(prs.slides), "rendered_pages": min(len(prs.slides), max_slides)},
        }

    # ---- archive (swiftpm/zip) → HTML ----
    def _preview_swiftpm(self, file_path: Path, html_mod) -> Dict[str, Any]:
        """解壓壓縮包，提取源碼並轉為 HTML 預覽。支持 .swiftpm 和 .zip。"""
        import zipfile

        if not zipfile.is_zipfile(file_path):
            raise ValidationError("此文件無法解壓（非有效 ZIP 格式）")

        max_total_chars = 100_000
        max_lines_per_file = 500
        total_chars = 0
        truncated = False
        file_blocks: List[str] = []
        code_exts = {".swift", ".py", ".js", ".ts", ".java", ".c", ".cpp", ".h", ".html", ".css", ".json"}

        with zipfile.ZipFile(file_path, "r") as zf:
            # 篩選源碼文件，排除系統文件
            swift_entries = sorted([
                n for n in zf.namelist()
                if any(n.endswith(ce) for ce in code_exts)
                and not n.startswith("__MACOSX/")
                and "/.build/" not in n
                and not n.split("/")[-1].startswith(".")
            ])

            for entry_name in swift_entries:
                try:
                    raw = zf.read(entry_name)
                    text = raw.decode("utf-8", errors="replace")
                except Exception:
                    continue

                # 清除 Swift Playgrounds 佔位標記
                text = re.sub(r'/\*@[A-Z_]+@\*/', '', text)

                lines = text.split("\n")
                if len(lines) > max_lines_per_file:
                    lines = lines[:max_lines_per_file]
                    truncated = True
                code_text = "\n".join(lines)

                if total_chars + len(code_text) > max_total_chars:
                    remaining = max_total_chars - total_chars
                    if remaining > 0:
                        code_text = code_text[:remaining]
                    truncated = True

                escaped = html_mod.escape(code_text)
                # 取最後兩層路徑作為顯示名
                display_name = "/".join(entry_name.rsplit("/", 2)[-2:]) if "/" in entry_name else entry_name
                file_blocks.append(
                    f'<div class="swiftpm-file">'
                    f'<div class="swiftpm-file-name">{html_mod.escape(display_name)}</div>'
                    f'<pre><code>{escaped}</code></pre>'
                    f'</div>'
                )

                total_chars += len(code_text)
                if total_chars >= max_total_chars:
                    break

        if not file_blocks:
            raise ValidationError("此壓縮包中未找到源碼文件")

        html_content = f'<div class="doc-preview swiftpm-preview">{"".join(file_blocks)}</div>'
        return {
            "html": html_content,
            "truncated": truncated,
            "meta": {"files_count": len(file_blocks)},
        }

    # ================================================================
    # Swift 運行
    # ================================================================

    @staticmethod
    def _find_swift_env():
        """查找 swiftc 可執行檔路徑及所需的環境變數"""
        import shutil
        import glob as _glob

        swift_base = os.path.join(
            os.environ.get("LOCALAPPDATA", ""), "Programs", "Swift"
        )

        # 1. 嘗試系統 PATH
        swiftc = shutil.which("swiftc")
        if swiftc:
            return swiftc, None

        # 2. Windows 默認安裝位置
        tc_pattern = os.path.join(swift_base, "Toolchains", "*", "usr", "bin", "swiftc.exe")
        matches = sorted(_glob.glob(tc_pattern), reverse=True)
        if not matches:
            return "swiftc", None

        swiftc_exe = matches[0]
        toolchain_bin = os.path.dirname(swiftc_exe)

        # 收集 DLL 路徑
        extra_paths = [toolchain_bin]
        rt_pattern = os.path.join(swift_base, "Runtimes", "*", "usr", "bin")
        for p in sorted(_glob.glob(rt_pattern), reverse=True):
            extra_paths.append(p)

        env = os.environ.copy()
        env["PATH"] = ";".join(extra_paths) + ";" + env.get("PATH", "")

        # SDKROOT
        sdk_pattern = os.path.join(swift_base, "Platforms", "*", "Windows.platform", "Developer", "SDKs", "Windows.sdk")
        sdk_matches = sorted(_glob.glob(sdk_pattern), reverse=True)
        if sdk_matches:
            env["SDKROOT"] = sdk_matches[0]

        # Windows SDK (UCRT) include / lib paths
        winsdk_base = os.path.join(os.environ.get("ProgramFiles(x86)", r"C:\Program Files (x86)"), "Windows Kits", "10")
        sdk_versions = sorted(_glob.glob(os.path.join(winsdk_base, "Include", "*")), reverse=True)
        if sdk_versions:
            sv = os.path.basename(sdk_versions[0])
            inc = os.path.join(winsdk_base, "Include", sv)
            lib = os.path.join(winsdk_base, "Lib", sv)
            env["INCLUDE"] = ";".join([
                os.path.join(inc, "ucrt"), os.path.join(inc, "um"), os.path.join(inc, "shared"),
            ]) + ";" + env.get("INCLUDE", "")
            env["LIB"] = ";".join([
                os.path.join(lib, "ucrt", "x64"), os.path.join(lib, "um", "x64"),
            ]) + ";" + env.get("LIB", "")

        # VC Tools include / lib
        vc_pattern = os.path.join(os.environ.get("ProgramFiles(x86)", r"C:\Program Files (x86)"),
                                  "Microsoft Visual Studio", "2022", "BuildTools", "VC", "Tools", "MSVC", "*")
        vc_matches = sorted(_glob.glob(vc_pattern), reverse=True)
        if vc_matches:
            vc = vc_matches[0]
            env["INCLUDE"] = os.path.join(vc, "include") + ";" + env.get("INCLUDE", "")
            env["LIB"] = os.path.join(vc, "lib", "x64") + ";" + env.get("LIB", "")

        return swiftc_exe, env

    def run_swift_code(self, code: str) -> Dict[str, Any]:
        """
        運行 Swift 代碼

        使用 swiftc 編譯後執行。
        """
        import subprocess
        import tempfile

        tmp_file = None
        exe_path = None
        try:
            swiftc_path, env = self._find_swift_env()

            # 寫入臨時文件
            tmp_file = tempfile.NamedTemporaryFile(
                suffix=".swift", mode="w", delete=False, encoding="utf-8"
            )
            tmp_file.write(code)
            tmp_file.close()
            exe_path = tmp_file.name.replace(".swift", ".exe")

            # 編譯
            compile_result = subprocess.run(
                [swiftc_path, "-o", exe_path, tmp_file.name],
                capture_output=True, text=True, timeout=30, env=env,
            )
            if compile_result.returncode != 0:
                return {
                    "success": False,
                    "stdout": "",
                    "stderr": compile_result.stderr,
                    "return_code": compile_result.returncode,
                }

            # 執行
            result = subprocess.run(
                [exe_path],
                capture_output=True, text=True, timeout=10, env=env,
            )

            return {
                "success": result.returncode == 0,
                "stdout": result.stdout,
                "stderr": result.stderr,
                "return_code": result.returncode,
            }

        except subprocess.TimeoutExpired:
            return {
                "success": False,
                "stdout": "",
                "stderr": "執行超時 (10秒限制)",
                "return_code": -1,
            }
        except FileNotFoundError:
            return {
                "success": False,
                "stdout": "",
                "stderr": "Swift 編譯器未找到，請確認已安裝 Swift",
                "return_code": -1,
            }
        except Exception as e:
            return {
                "success": False,
                "stdout": "",
                "stderr": str(e),
                "return_code": -1,
            }
        finally:
            # 清理臨時文件
            if tmp_file and os.path.exists(tmp_file.name):
                os.unlink(tmp_file.name)
            if exe_path and os.path.exists(exe_path):
                os.unlink(exe_path)
            # 清理編譯產生的 .lib / .exp 文件
            if exe_path:
                for ext in (".lib", ".exp"):
                    p = exe_path.replace(".exe", ext)
                    if os.path.exists(p):
                        os.unlink(p)

    # ================================================================
    # 作業 AI 問答
    # ================================================================

    def build_assignment_context(self, assignment_id: int, user_id: int) -> str:
        """
        構建作業 AI 問答的上下文字符串。
        驗證學生擁有此提交且已批改，然後組裝作業信息、提交內容、批改結果。

        Returns:
            str: 結構化上下文字符串（含系統指令）
        """
        assignment = self._get_assignment_or_raise(assignment_id)
        self._deserialize_json_fields(assignment)

        # 查詢學生的提交
        submission = self._submission_repo.find_one(
            where="assignment_id = %s AND student_id = %s",
            params=(assignment_id, user_id),
        )
        if not submission:
            raise ValidationError("你尚未提交此作業")
        if submission.get("status") != "graded":
            raise ValidationError("此作業尚未批改完成")

        # 評分標準
        rubric_items = self._rubric_repo.find_by_assignment(assignment_id)
        for item in rubric_items:
            self._deserialize_json_fields(item)

        # 逐項得分
        rubric_scores = self._score_repo.find_by_submission(submission["id"])
        score_map = {}
        level_map = {}
        for s in rubric_scores:
            score_map[s.get("rubric_item_id")] = s.get("points")
            if s.get("selected_level"):
                level_map[s.get("rubric_item_id")] = s["selected_level"]

        # 提交文件名
        files = self._file_repo.find_by_submission(submission["id"])
        file_names = ", ".join(f.get("original_name", "") for f in files) if files else "（無文件）"

        # 評分方式顯示名
        rubric_type = assignment.get("rubric_type", "points")
        type_display = {
            "points": "計分制",
            "analytic_levels": "分級量規",
            "weighted_pct": "權重百分比制",
            "checklist": "通過/不通過清單",
            "competency": "能力等級制",
            "dse_criterion": "DSE 標準量規",
            "holistic": "整體評分",
        }.get(rubric_type, rubric_type)

        # 組裝評分項信息
        rubric_lines = []
        for item in rubric_items:
            max_pts = item.get("max_points", 0)
            rubric_lines.append(f"  - {item.get('title', '')}（滿分 {max_pts}）")
        rubric_str = "\n".join(rubric_lines) if rubric_lines else "  （無評分項）"

        # 組裝各項得分
        score_lines = []
        for item in rubric_items:
            item_id = item.get("id")
            pts = score_map.get(item_id, "—")
            max_pts = item.get("max_points", 0)
            level = level_map.get(item_id, "")
            level_str = f" ({level})" if level else ""
            score_lines.append(f"  - {item.get('title', '')}：{pts}/{max_pts}{level_str}")
        scores_str = "\n".join(score_lines) if score_lines else "  （無逐項得分）"

        feedback = submission.get("feedback") or "（教師未留下評語）"
        total_score = submission.get("score", "—")
        max_score = assignment.get("max_score", "—")

        context = f"""【你的身份】
你是一位耐心的 AI 學習助教。學生正在查看一份已批改的作業，想要了解自己的表現和相關知識。

【重要規則】
1. 只能依據下方【批改結果】中已有的資訊來回答關於得分的問題
2. 如果教師沒有針對某項給出具體扣分原因，你必須明確告知學生「現有批改資料中未說明具體扣分原因，建議向老師確認」
3. 你可以根據作業要求和學生提交內容，給出「可能的改進方向」，但必須標明這是你的建議而非教師的評價
4. 絕對不能把推測說成事實
5. 回答要基於知識庫中的學科知識，幫助學生理解相關概念

【作業信息】
標題：{assignment.get('title', '')}
要求：{assignment.get('description', '（無描述）')}
評分方式：{type_display}
評分項：
{rubric_str}

【學生提交】
提交內容：{submission.get('content', '（無備註）')}
提交文件：{file_names}

【批改結果】
總分：{total_score} / {max_score}
各項得分：
{scores_str}
教師評語：{feedback}

【學生提問】"""
        return context

    def validate_ai_conversation(
        self, conversation_id: str, assignment_id: int, username: str
    ) -> bool:
        """
        驗證 conversation 屬於該學生且對應同一份作業。
        防止學生手動修改 conversation_id 來跨作業串話。
        """
        if not self._conversation_repo:
            logger.error("validate_ai_conversation: conversation_repo 未注入")
            return False

        conv = self._conversation_repo.find_one(
            "conversation_id = %s AND (is_deleted = 0 OR is_deleted IS NULL)",
            (conversation_id,),
        )
        if not conv:
            return False
        return (
            conv.get("username") == username
            and conv.get("assignment_id") == assignment_id
        )

    # ================================================================
    # Form 作業
    # ================================================================

    def _create_form_questions(self, assignment_id: int, questions: List[Dict]) -> List[int]:
        """批量創建 Form 題目 + MC 選項"""
        if not self._question_repo or not self._question_option_repo:
            raise ValidationError("Form 功能未初始化")

        question_ids = self._question_repo.batch_insert(assignment_id, questions)
        for qid, q in zip(question_ids, questions):
            if q.get("question_type") == "mc" and q.get("options"):
                opts = q["options"]
                if isinstance(opts, list) and opts and isinstance(opts[0], dict):
                    self._question_option_repo.batch_insert(qid, opts)
                elif isinstance(opts, list) and opts and hasattr(opts[0], "dict"):
                    self._question_option_repo.batch_insert(
                        qid, [o.dict() for o in opts]
                    )
        return question_ids

    def get_form_questions(
        self, assignment_id: int, include_answers: bool = True
    ) -> List[Dict[str, Any]]:
        """獲取 Form 題目列表。include_answers=False 時隱藏正確答案和參考答案（學生端）"""
        if not self._question_repo:
            return []

        questions = self._question_repo.find_by_assignment(assignment_id)
        if not questions:
            return []

        # 批量獲取選項（若 option repo 可用）
        options_by_q = {}
        if self._question_option_repo:
            q_ids = [q["id"] for q in questions]
            all_options = self._question_option_repo.find_by_questions(q_ids)
            for opt in all_options:
                options_by_q.setdefault(opt["question_id"], []).append(opt)

        result = []
        for q in questions:
            # 兼容 exam-paper 和 form 兩種 schema
            pts = q.get("max_points") or q.get("points")
            # 解析 metadata
            raw_md = q.get("metadata")
            md = {}
            if isinstance(raw_md, str):
                try:
                    md = json.loads(raw_md)
                except (json.JSONDecodeError, TypeError):
                    md = {}
            elif isinstance(raw_md, dict):
                md = raw_md

            # 選項: 優先 option_repo，否則 metadata.options
            opts = options_by_q.get(q["id"], [])
            if not opts and md.get("options"):
                raw_opts = md["options"]
                if isinstance(raw_opts, list) and raw_opts:
                    if isinstance(raw_opts[0], str):
                        keys = "ABCDEFGHIJKLMNOP"
                        opts = [{"option_key": keys[j] if j < len(keys) else str(j+1),
                                 "option_text": o} for j, o in enumerate(raw_opts)]
                    elif isinstance(raw_opts[0], dict):
                        opts = raw_opts

            item = {
                "id": q["id"],
                "question_order": q.get("question_order", 0),
                "question_number": q.get("question_number", ""),
                "question_type": q.get("question_type", "open"),
                "question_text": q.get("question_text", ""),
                "max_points": float(pts) if pts is not None else 0,
                "grading_notes": q.get("grading_notes") or "",
                "options": opts,
                "metadata": md,
            }
            if include_answers:
                item["correct_answer"] = q.get("correct_answer") or q.get("answer_text") or ""
                item["reference_answer"] = q.get("reference_answer") or ""
            result.append(item)

        return result

    async def submit_form_answers(
        self,
        assignment_id: int,
        student: Dict,
        answers: List[Dict],
        files_by_question: Optional[Dict[int, List[UploadFile]]] = None,
    ) -> Dict[str, Any]:
        """學生提交 Form 作業（事務性）"""
        assignment = self._get_assignment_or_raise(assignment_id)

        # 校驗
        if (assignment.get("assignment_type") or AssignmentType.FILE_UPLOAD) != AssignmentType.FORM:
            raise ValidationError("此作業不是 Form 類型")
        if assignment["status"] != "published":
            raise AssignmentNotPublishedError()

        # 截止時間 server-side 強制
        if assignment.get("deadline"):
            deadline = assignment["deadline"]
            if isinstance(deadline, str):
                deadline = datetime.fromisoformat(deadline)
            if datetime.now() > deadline and not assignment.get("allow_late"):
                raise DeadlinePassedError()

        is_late = False
        if assignment.get("deadline"):
            deadline = assignment["deadline"]
            if isinstance(deadline, str):
                deadline = datetime.fromisoformat(deadline)
            is_late = datetime.now() > deadline

        # 重複提交檢查
        student_id = student.get("id") or student.get("user_id")
        existing = self._submission_repo.find_by_assignment_student(assignment_id, student_id)
        if existing:
            raise ValidationError("已提交過此作業，不可重複提交")

        # 獲取題目
        if not self._question_repo:
            raise ValidationError("Form 功能未初始化")
        questions = self._question_repo.find_by_assignment(assignment_id)
        if not questions:
            raise ValidationError("此作業沒有題目")

        q_map = {q["id"]: q for q in questions}
        q_ids_set = set(q_map.keys())

        # 驗證 answers 覆蓋所有題目
        submitted_q_ids = set(a["question_id"] for a in answers)
        if submitted_q_ids != q_ids_set:
            missing = q_ids_set - submitted_q_ids
            extra = submitted_q_ids - q_ids_set
            msg = []
            if missing:
                msg.append(f"缺少題目: {missing}")
            if extra:
                msg.append(f"多餘題目: {extra}")
            raise ValidationError("; ".join(msg))

        # 創建 submission
        submission_id = self._submission_repo.insert_get_id({
            "assignment_id": assignment_id,
            "student_id": student_id,
            "student_name": student.get("display_name") or student.get("username", ""),
            "username": student.get("username", ""),
            "class_name": student.get("class_name", ""),
            "content": "",
            "status": "submitted",
            "is_late": is_late,
            "submitted_at": datetime.now(),
        })

        # 寫入作答 + MC 自動評分
        mc_total = 0.0
        answer_records = []
        for ans in answers:
            q = q_map[ans["question_id"]]
            record = {
                "question_id": ans["question_id"],
                "answer_text": ans.get("answer_text", ""),
            }

            if q["question_type"] == "mc":
                correct = (q.get("correct_answer") or "").strip().upper()
                student_ans = (ans.get("answer_text") or "").strip().upper()
                is_correct = student_ans == correct if correct else False
                raw_mp = q.get("max_points") or q.get("points") or 0
                pts = float(raw_mp) if is_correct else 0.0
                record["is_correct"] = is_correct
                record["points"] = pts
                record["score_source"] = ScoreSource.AUTO
                mc_total += pts
            # 文字題：points=NULL, score_source=NULL

            answer_records.append(record)

        if not self._answer_repo:
            raise ValidationError("Form 功能未初始化")
        answer_ids = self._answer_repo.batch_insert(submission_id, answer_records)

        # 保存文字題上傳文件
        if files_by_question and self._answer_file_repo:
            # 建立 question_id → answer_id 映射
            q_to_answer = {}
            for aid, ans in zip(answer_ids, answers):
                q_to_answer[ans["question_id"]] = aid

            for q_id, file_list in files_by_question.items():
                if q_id not in q_to_answer:
                    continue
                answer_id = q_to_answer[q_id]
                for f in file_list:
                    await self._save_answer_file(answer_id, f)

        # 回寫 submission 總分（目前只有 MC 部分）
        self._submission_repo.update(
            data={"score": mc_total},
            where="id = %s",
            params=(submission_id,),
        )

        logger.info("學生 %s 提交了 Form 作業 #%d (submission_id=%d)",
                     student.get("username"), assignment_id, submission_id)
        return self.get_submission_detail(submission_id)

    async def submit_exam_answers(
        self,
        assignment_id: int,
        student: Dict,
        answers: List[Dict],
    ) -> Dict[str, Any]:
        """學生提交 Exam 類型作業 (試卷題目作答)"""
        assignment = self._get_assignment_or_raise(assignment_id)

        if (assignment.get("assignment_type") or AssignmentType.FILE_UPLOAD) != AssignmentType.EXAM:
            raise ValidationError("此作業不是 Exam 類型")
        if assignment["status"] != "published":
            raise AssignmentNotPublishedError()

        # 截止時間
        if assignment.get("deadline"):
            deadline = assignment["deadline"]
            if isinstance(deadline, str):
                deadline = datetime.fromisoformat(deadline)
            if datetime.now() > deadline and not assignment.get("allow_late"):
                raise DeadlinePassedError()

        is_late = False
        if assignment.get("deadline"):
            deadline = assignment["deadline"]
            if isinstance(deadline, str):
                deadline = datetime.fromisoformat(deadline)
            is_late = datetime.now() > deadline

        student_id = student.get("id") or student.get("user_id")
        existing = self._submission_repo.find_by_assignment_student(assignment_id, student_id)

        # 獲取試卷題目
        questions = self.get_assignment_questions(assignment_id)
        if not questions:
            raise ValidationError("此作業沒有題目")

        q_map = {q["id"]: q for q in questions}
        # 排除 passage 類型 (不需作答)
        answerable_ids = {
            q["id"] for q in questions
            if q.get("question_type") != "passage"
        }

        if existing:
            # 更新現有提交 (截止前可重複提交)
            submission_id = existing["id"]
            self._submission_repo.update(
                data={
                    "status": "submitted",
                    "is_late": is_late,
                    "score": None,
                    "feedback": None,
                    "graded_by": None,
                    "graded_at": None,
                    "submitted_at": datetime.now(),
                    "updated_at": datetime.now(),
                },
                where="id = %s",
                params=(submission_id,),
            )
            # 刪除舊作答
            if self._answer_repo:
                self._answer_repo.delete_by_submission(submission_id)
        else:
            # 創建新提交
            submission_id = self._submission_repo.insert_get_id({
                "assignment_id": assignment_id,
                "student_id": student_id,
                "student_name": student.get("display_name") or student.get("username", ""),
                "username": student.get("username", ""),
                "class_name": student.get("class_name", ""),
                "content": "",
                "status": "submitted",
                "is_late": is_late,
                "submitted_at": datetime.now(),
            })

        # 寫入作答記錄
        answer_records = []
        for ans in answers:
            q_id = ans.get("question_id")
            if q_id not in answerable_ids:
                continue
            record = {
                "question_id": q_id,
                "answer_text": ans.get("answer_text", ""),
            }
            answer_records.append(record)

        if self._answer_repo and answer_records:
            self._answer_repo.batch_insert(submission_id, answer_records)

        logger.info("學生 %s 提交了 Exam 作業 #%d (submission_id=%d)",
                     student.get("username"), assignment_id, submission_id)
        return self.get_submission_detail(submission_id)

    async def _save_answer_file(self, answer_id: int, file: UploadFile) -> Dict[str, Any]:
        """保存作答附件"""
        original_name = file.filename or "unnamed"
        ext = Path(original_name).suffix.lower()

        file_type = EXTENSION_TYPE_MAP.get(ext, "")
        content = await file.read()
        file_size = len(content)

        if file_size > MAX_FILE_SIZE:
            raise FileTooLargeError(MAX_FILE_SIZE // 1024 // 1024)

        stored_name = f"{uuid.uuid4().hex}{ext}"
        file_path = UPLOAD_DIR / stored_name

        with open(file_path, "wb") as fh:
            fh.write(content)

        self._answer_file_repo.insert({
            "answer_id": answer_id,
            "original_name": original_name,
            "stored_name": stored_name,
            "file_path": f"uploads/assignments/{stored_name}",
            "file_size": file_size,
            "file_type": file_type,
            "mime_type": file.content_type or "",
        })

        return {"original_name": original_name, "stored_name": stored_name}

    async def ai_grade_form_submission(self, submission_id: int, extra_prompt: str = "") -> Dict[str, Any]:
        """AI 批改 Form 提交中的文字題"""
        if not self._ask_ai_func:
            raise ValidationError("AI 功能未初始化")
        if not self._answer_repo:
            raise ValidationError("Form 功能未初始化")

        submission = self._submission_repo.find_by_id(submission_id)
        if not submission:
            raise SubmissionNotFoundError(submission_id)

        assignment = self._get_assignment_or_raise(submission["assignment_id"])
        questions = self._question_repo.find_by_assignment(assignment["id"]) if self._question_repo else []
        q_map = {q["id"]: q for q in questions}

        # 找出需要 AI 批改的答案
        ungraded = self._answer_repo.find_ungraded_text(submission_id)
        if not ungraded:
            return {"graded_count": 0, "message": "沒有需要批改的文字題"}

        graded_count = 0
        for ans in ungraded:
            q = q_map.get(ans["question_id"])
            if not q or q["question_type"] == "mc":
                continue

            # 兼容 exam-paper (max_points) 和 form (points) 兩種 schema
            raw_pts = q.get("max_points") or q.get("points")
            max_pts = float(raw_pts) if raw_pts is not None else 0
            reference = (q.get("reference_answer") or q.get("answer_text") or "").strip()
            grading_notes = (q.get("grading_notes") or "").strip()

            # 教師額外批改指示
            teacher_note = ""
            if extra_prompt and extra_prompt.strip():
                teacher_note = f"教師批改指示：{extra_prompt.strip()}\n"

            if reference:
                prompt = (
                    f"你是一位嚴謹的教師，正在根據參考答案批改學生的答案。\n\n"
                    f"題目：{q['question_text']}\n"
                    f"題型：{'短答' if q['question_type'] == 'short_answer' else '長答'}\n"
                    f"滿分：{max_pts}\n"
                    f"老師參考答案：{reference}\n"
                )
                if grading_notes:
                    prompt += f"批改注意事項：{grading_notes}\n"
                if teacher_note:
                    prompt += teacher_note
                prompt += (
                    f"\n學生答案：{ans.get('answer_text') or '（空白）'}\n\n"
                    f"請嚴格按照參考答案的要點來評分。\n"
                    f'請用 JSON 格式回覆：{{"points": <0到{max_pts}的數字>, "feedback": "<評語>"}}'
                )
            else:
                prompt = (
                    f"你是一位專業教師，正在批改學生的答案。\n\n"
                    f"題目：{q['question_text']}\n"
                    f"題型：{'短答' if q['question_type'] == 'short_answer' else '長答'}\n"
                    f"滿分：{max_pts}\n"
                )
                if grading_notes:
                    prompt += f"批改注意事項：{grading_notes}\n"
                if teacher_note:
                    prompt += teacher_note
                prompt += (
                    f"\n學生答案：{ans.get('answer_text') or '（空白）'}\n\n"
                    f"請根據答案的正確性和完整性來評分。\n"
                    f'請用 JSON 格式回覆：{{"points": <0到{max_pts}的數字>, "feedback": "<評語>"}}'
                )

            try:
                # 與 rubric 批改保持一致的調用方式
                ai_response, _ = self._ask_ai_func(
                    question=prompt,
                    subject_code=assignment.get("subject", "general"),
                    use_api=False,
                    conversation_history=[],
                )

                ai_result = self._parse_ai_form_response(str(ai_response), max_pts)
                ai_pts = ai_result.get("points", 0)
                ai_feedback = ai_result.get("feedback", "")

                update_data = {
                    "ai_points": ai_pts,
                    "ai_feedback": ai_feedback,
                    "graded_at": datetime.now(),
                }
                if not ans.get("reviewed_at"):
                    update_data["points"] = ai_pts
                    update_data["score_source"] = ScoreSource.AI

                self._answer_repo.update_score(ans["id"], update_data)
                graded_count += 1

            except Exception as e:
                logger.error("AI 批改 answer #%d 失敗: %s", ans["id"], e)

        # 回寫總分
        self._recalculate_form_score(submission_id)

        return {"graded_count": graded_count, "message": f"已批改 {graded_count} 道文字題"}

    def _parse_ai_form_response(self, response: str, max_points: float) -> Dict[str, Any]:
        """解析 AI 批改回覆為 {points, feedback}"""
        if not response:
            return {"points": 0, "feedback": "AI 未返回結果"}

        try:
            result = json.loads(response)
            if isinstance(result, dict) and "points" in result:
                pts = min(max(float(result["points"]), 0), max_points)
                return {"points": pts, "feedback": result.get("feedback", "")}
        except (json.JSONDecodeError, ValueError):
            pass

        json_match = re.search(r'\{[^{}]*"points"[^{}]*\}', response, re.DOTALL)
        if json_match:
            try:
                result = json.loads(json_match.group())
                pts = min(max(float(result.get("points", 0)), 0), max_points)
                return {"points": pts, "feedback": result.get("feedback", "")}
            except (json.JSONDecodeError, ValueError):
                pass

        return {"points": 0, "feedback": response[:500]}

    def teacher_grade_form_answer(
        self,
        answer_id: int,
        teacher_id: int,
        points: float,
        feedback: str = "",
    ) -> Dict[str, Any]:
        """教師手動批改 Form 單題"""
        if not self._answer_repo:
            raise ValidationError("Form 功能未初始化")

        answer = self._answer_repo.find_by_id(answer_id)
        if not answer:
            raise NotFoundError("找不到作答記錄")

        if self._question_repo:
            question = self._question_repo.find_by_id(answer["question_id"])
            max_pts = question.get("max_points") or question.get("points") if question else None
            if max_pts is not None and points > float(max_pts):
                raise ValidationError(f"分數不能超過滿分 {max_pts}")

        self._answer_repo.update_score(answer_id, {
            "points": points,
            "teacher_feedback": feedback,
            "score_source": ScoreSource.TEACHER,
            "reviewed_at": datetime.now(),
            "graded_at": datetime.now(),
        })

        self._recalculate_form_score(answer["submission_id"])

        logger.info("教師手動批改 answer #%d: %.1f 分", answer_id, points)
        return self.get_submission_detail(answer["submission_id"])

    def _recalculate_form_score(self, submission_id: int) -> None:
        """重新計算 Form 提交總分"""
        if not self._answer_repo:
            return

        answers = self._answer_repo.find_by_submission(submission_id)
        total = sum(
            float(a["points"])
            for a in answers
            if a.get("points") is not None
        )

        self._submission_repo.update(
            data={"score": total},
            where="id = %s",
            params=(submission_id,),
        )

    # ================================================================
    # 輔助方法
    # ================================================================

    def _get_assignment_or_raise(self, assignment_id: int) -> Dict[str, Any]:
        """獲取作業，不存在或已刪除則拋出異常"""
        assignment = self._assignment_repo.find_one(
            where="id = %s AND is_deleted = 0",
            params=(assignment_id,),
        )
        if not assignment:
            raise AssignmentNotFoundError(assignment_id)
        return assignment

    # ================================================================
    # 試卷上傳 + OCR 識別
    # ================================================================

    def create_upload_batch(
        self,
        batch_id: str,
        subject: str,
        created_by: int,
        total_files: int,
    ) -> None:
        """創建上傳批次記錄"""
        self._batch_repo.create_batch(batch_id, subject, created_by, total_files)

    def create_upload_file(self, data: Dict[str, Any]) -> None:
        """創建上傳文件記錄"""
        self._upload_file_repo.create_file(data)

    def get_batch(self, batch_id: str) -> Optional[Dict[str, Any]]:
        """獲取批次記錄"""
        return self._batch_repo.find_by_batch_id(batch_id)

    def get_batch_files(self, batch_id: str) -> List[Dict[str, Any]]:
        """獲取批次的所有文件"""
        return self._upload_file_repo.find_by_batch(batch_id)

    def update_file_ocr(
        self,
        file_id: int,
        status: str,
        result: Optional[str] = None,
        error: Optional[str] = None,
    ) -> None:
        """更新文件 OCR 狀態"""
        self._upload_file_repo.update_ocr_status(file_id, status, result, error)

    def update_batch_status(self, batch_id: str, status: str, **kwargs) -> None:
        """更新批次聚合狀態"""
        self._batch_repo.update_status(batch_id, status, **kwargs)

    def refresh_batch_status(self, batch_id: str) -> Dict[str, Any]:
        """根據文件狀態刷新批次聚合狀態"""
        files = self._upload_file_repo.find_by_batch(batch_id)
        total = len(files)
        completed = sum(1 for f in files if f["ocr_status"] == "completed")
        failed = sum(1 for f in files if f["ocr_status"] == "failed")

        # 計算題目數和低置信度數
        total_questions = 0
        low_confidence_count = 0
        for f in files:
            if f["ocr_status"] == "completed" and f.get("ocr_result"):
                try:
                    questions = json.loads(f["ocr_result"]) if isinstance(f["ocr_result"], str) else f["ocr_result"]
                    if isinstance(questions, list):
                        total_questions += len(questions)
                        low_confidence_count += sum(
                            1 for q in questions
                            if isinstance(q, dict) and q.get("confidence", 1.0) < 0.6
                        )
                except (json.JSONDecodeError, TypeError):
                    pass

        if completed == total:
            status = "completed"
        elif failed == total:
            status = "failed"
        elif completed + failed == total and failed > 0:
            status = "partial_failed"
        else:
            status = "processing"

        self._batch_repo.update_status(
            batch_id, status,
            completed_files=completed,
            failed_files=failed,
            total_questions=total_questions,
            low_confidence_count=low_confidence_count,
        )
        return {
            "status": status,
            "completed_files": completed,
            "failed_files": failed,
            "total_questions": total_questions,
            "low_confidence_count": low_confidence_count,
        }

    def get_batch_status(self, batch_id: str) -> Optional[Dict[str, Any]]:
        """獲取批次聚合狀態 + 文件詳情 + 合併題目"""
        batch = self._batch_repo.find_by_batch_id(batch_id)
        if not batch:
            return None

        files = self._upload_file_repo.find_by_batch(batch_id)
        file_details = []
        merged_questions = []

        for f in files:
            detail = {
                "id": f["id"],
                "filename": f["original_filename"],
                "file_type": f["file_type"],
                "status": f["ocr_status"],
                "total_pages": f.get("total_pages", 1),
                "error": f.get("error_message"),
                "question_count": 0,
            }
            if f["ocr_status"] == "completed" and f.get("ocr_result"):
                try:
                    questions = json.loads(f["ocr_result"]) if isinstance(f["ocr_result"], str) else f["ocr_result"]
                    if isinstance(questions, list):
                        detail["question_count"] = len(questions)
                        merged_questions.extend(questions)
                except (json.JSONDecodeError, TypeError):
                    pass
            file_details.append(detail)

        result = {
            "batch_id": batch["batch_id"],
            "subject": batch["subject"],
            "status": batch["status"],
            "total_files": batch["total_files"],
            "completed_files": batch.get("completed_files", 0),
            "failed_files": batch.get("failed_files", 0),
            "total_questions": batch.get("total_questions", 0),
            "low_confidence_count": batch.get("low_confidence_count", 0),
            "files": file_details,
        }

        # completed 或 partial_failed 都返回已完成文件的題目
        if batch["status"] in ("completed", "partial_failed"):
            result["questions"] = merged_questions
            if batch["status"] == "partial_failed":
                result["warning"] = f"部分文件識別失敗 ({batch.get('failed_files', 0)} 個)"
        else:
            result["questions"] = []

        return result

    # ================================================================
    # 作業題目 CRUD
    # ================================================================

    def get_assignment_questions(self, assignment_id: int) -> List[Dict[str, Any]]:
        """獲取作業題目列表"""
        self._get_assignment_or_raise(assignment_id)
        questions = self._question_repo.find_by_assignment(assignment_id)
        for q in questions:
            self._deserialize_json_fields(q)
        return questions

    def save_assignment_questions(
        self,
        assignment_id: int,
        questions: List[Dict[str, Any]],
    ) -> List[Dict[str, Any]]:
        """
        保存/替換作業題目 (事務化)。

        規則:
        - 先校驗再刪除
        - draft / 無 submission → 允許
        - 已有 submission → 禁止
        """
        assignment = self._get_assignment_or_raise(assignment_id)

        # 校驗: 如果已有 submission，禁止修改
        submission_count = self._submission_repo.count(
            where="assignment_id = %s", params=(assignment_id,)
        )
        if submission_count > 0:
            raise ValidationError(
                "作業已有學生提交，無法修改題目",
                field="questions",
            )

        # 校驗 payload
        if not questions:
            raise ValidationError("題目列表不能為空", field="questions")

        for i, q in enumerate(questions):
            q_type = q.get("question_type", "open")

            # passage 類型不需要作答，但仍需 question_text
            if not q.get("question_text", "").strip():
                raise ValidationError(
                    f"第 {i + 1} 題的題目文字不能為空",
                    field="questions",
                )

            # fill_blank: 校驗 + 規範化 blanks
            if q_type == "fill_blank":
                md = q.get("metadata") or {}
                blanks_raw = md.get("blanks", [])
                if not isinstance(blanks_raw, list) or not blanks_raw:
                    raise ValidationError(
                        f"第 {i + 1} 題 (填空題) 至少需要一個填空項",
                        field="questions",
                    )
                has_template = bool(md.get("template_text", "").strip())
                validated_blanks = []
                seen_ids = set()
                for bi, b in enumerate(blanks_raw):
                    if not isinstance(b, dict):
                        continue
                    blank_id = str(b.get("id") or f"b{bi + 1}").strip()
                    label = str(b.get("label", "")).strip()
                    # 模板模式下 label 可為空；非模板模式下仍需 label
                    if not label and not has_template:
                        raise ValidationError(
                            f"第 {i + 1} 題的第 {bi + 1} 個填空項缺少標籤",
                            field="questions",
                        )
                    pts_raw = b.get("points")
                    try:
                        pts = float(pts_raw) if pts_raw is not None else 0.0
                    except (ValueError, TypeError):
                        pts = 0.0
                    if blank_id in seen_ids:
                        blank_id = f"b{bi + 1}"
                    seen_ids.add(blank_id)
                    # input_type 校驗
                    input_type = str(b.get("input_type", "short_text")).strip()
                    if input_type not in ("short_text", "long_text"):
                        input_type = "short_text"
                    validated_blanks.append({
                        "id": blank_id,
                        "label": label,
                        "points": round(pts, 1),
                        "answer": str(b.get("answer", "")).strip(),
                        "input_type": input_type,
                    })
                # auto-sum points (後端覆蓋，不信任前端)
                q["points"] = sum(b["points"] for b in validated_blanks)
                md["blanks"] = validated_blanks
                # blank_mode 校驗
                blank_mode = md.get("blank_mode", "inline")
                if blank_mode not in ("inline", "section", "mixed"):
                    blank_mode = "inline"
                md["blank_mode"] = blank_mode
                q["metadata"] = md
            else:
                # 非 fill_blank: 清除殘留 blanks 防髒數據
                md = q.get("metadata") or {}
                if isinstance(md, dict):
                    md.pop("blanks", None)
                    md.pop("blank_mode", None)
                    q["metadata"] = md if md else None

        # 事務化: delete + insert
        with self._question_repo.transaction() as conn:
            cursor = conn.cursor()
            cursor.execute(
                "DELETE FROM assignment_questions WHERE assignment_id = %s",
                (assignment_id,),
            )
            for i, q in enumerate(questions):
                md = q.get("metadata")
                metadata_str = None
                if md is not None:
                    metadata_str = json.dumps(md, ensure_ascii=False) if isinstance(md, (dict, list)) else md

                cursor.execute(
                    """INSERT INTO assignment_questions
                    (assignment_id, question_order, question_number, question_text,
                     answer_text, answer_source, points, question_type,
                     question_type_confidence, is_ai_extracted, source_batch_id,
                     source_page, ocr_confidence, metadata)
                    VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)""",
                    (
                        assignment_id, i,
                        q.get("question_number", ""),
                        q.get("question_text", ""),
                        q.get("answer_text", ""),
                        q.get("answer_source", "missing"),
                        q.get("points"),
                        q.get("question_type", "open"),
                        q.get("question_type_confidence"),
                        q.get("is_ai_extracted", True),
                        q.get("source_batch_id"),
                        q.get("source_page"),
                        q.get("ocr_confidence"),
                        metadata_str,
                    ),
                )
            # Exam 類型：更新 max_score = 非 passage 題目 points 之和
            total_points = sum(
                float(q.get("points") or 0)
                for q in questions
                if q.get("question_type") != "passage"
            )
            cursor.execute(
                "UPDATE assignments SET updated_at = %s, max_score = %s WHERE id = %s",
                (datetime.now(), total_points, assignment_id),
            )

        logger.info("作業 #%d 保存了 %d 道題目", assignment_id, len(questions))
        return self._question_repo.find_by_assignment(assignment_id)

    def get_assignment_questions_for_student(
        self, assignment_id: int
    ) -> List[Dict[str, Any]]:
        """獲取學生視角的題目 (過濾答案和內部字段, 保留 metadata 但去除答案)"""
        questions = self.get_assignment_questions(assignment_id)
        student_fields = [
            "id", "question_order", "question_number", "question_text",
            "points", "question_type", "metadata",
        ]
        result = []
        for q in questions:
            sq = {k: q.get(k) for k in student_fields if k in q}
            # fill_blank: 保留 blanks 結構但去除答案
            md = sq.get("metadata")
            if isinstance(md, dict) and md.get("blanks"):
                sq["metadata"] = {
                    **md,
                    "blanks": [
                        {k: v for k, v in b.items() if k != "answer"}
                        for b in md["blanks"]
                    ],
                }
            result.append(sq)
        return result

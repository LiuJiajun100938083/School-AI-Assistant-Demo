# llm/rag/file_processor.py
"""
文件处理器模块
支持 TXT/MD/PDF/DOCX/PPTX 等多格式文本提取，带多编码检测、PDF 多方法 fallback 和 OCR 支持。
"""
import os
import logging
from pathlib import Path
from typing import Tuple, Dict
import tempfile

logger = logging.getLogger(__name__)


class FileProcessor:
    """文件处理器 - 处理各种格式的文档"""

    def __init__(self):
        self.supported_formats = {
            '.txt': self._process_text,
            '.md': self._process_text,
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.pptx': self._process_pptx
        }

        # 检查可选依赖
        self.pdf_available = self._check_pdf_support()
        self.docx_available = self._check_docx_support()
        self.pptx_available = self._check_pptx_support()

        logger.info("📄 文件处理器初始化完成")
        logger.info(f"  PDF支持: {'✅' if self.pdf_available else '❌'}")
        logger.info(f"  DOCX支持: {'✅' if self.docx_available else '❌'}")
        logger.info(f"  PPTX支持: {'✅' if self.pptx_available else '❌'}")

    def _check_pdf_support(self) -> bool:
        """检查PDF处理支持"""
        try:
            import PyPDF2
            return True
        except ImportError:
            try:
                import pdfplumber
                return True
            except ImportError:
                logger.warning("⚠️ PDF处理库未安装 (PyPDF2 或 pdfplumber)")
                return False

    def _check_docx_support(self) -> bool:
        """检查DOCX处理支持"""
        try:
            import docx
            return True
        except ImportError:
            logger.warning("⚠️ DOCX处理库未安装 (python-docx)")
            return False

    def _check_pptx_support(self) -> bool:
        """检查PPTX处理支持"""
        try:
            import pptx
            return True
        except ImportError:
            logger.warning("⚠️ PPTX处理库未安装 (python-pptx)")
            return False

    def is_supported_format(self, filename: str) -> bool:
        """检查是否支持的文件格式"""
        ext = Path(filename).suffix.lower()

        if ext not in self.supported_formats:
            return False

        if ext == '.pdf' and not self.pdf_available:
            return False
        if ext == '.docx' and not self.docx_available:
            return False
        if ext == '.pptx' and not self.pptx_available:
            return False

        return True

    def get_supported_formats(self) -> list:
        """获取支持的文件格式列表"""
        formats = []
        for ext in self.supported_formats:
            if ext in ['.txt', '.md']:
                formats.append(ext)
            elif ext == '.pdf' and self.pdf_available:
                formats.append(ext)
            elif ext == '.docx' and self.docx_available:
                formats.append(ext)
            elif ext == '.pptx' and self.pptx_available:
                formats.append(ext)
        return formats

    def process_file(self, file_path: str, original_filename: str,
                     user: str = "unknown") -> Tuple[bool, str, Dict]:
        """
        处理文件并提取文本内容

        Returns:
            (成功标志, 提取的文本内容, 文件信息字典)
        """
        try:
            ext = Path(original_filename).suffix.lower()

            if not self.is_supported_format(original_filename):
                return False, "", {"error": f"不支持的文件格式: {ext}"}

            processor = self.supported_formats.get(ext)
            if processor:
                content = processor(file_path)

                file_stat = os.stat(file_path)
                file_info = {
                    "filename": original_filename,
                    "size": file_stat.st_size,
                    "format": ext,
                    "processed_by": user,
                    "content_length": len(content)
                }

                return True, content, file_info
            else:
                return False, "", {"error": f"没有找到处理器: {ext}"}

        except Exception as e:
            logger.error(f"❌ 处理文件失败 {original_filename}: {e}")
            return False, "", {"error": str(e)}

    def _process_text(self, file_path: str) -> str:
        """处理文本文件（TXT, MD）- 带多编码检测"""
        try:
            encodings = ['utf-8', 'gbk', 'gb2312', 'big5', 'utf-16']

            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                        logger.info(f"✅ 使用 {encoding} 编码成功读取文本文件")
                        return content
                except UnicodeDecodeError:
                    continue

            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
                logger.warning("⚠️ 使用UTF-8编码（忽略错误）读取文本文件")
                return content

        except Exception as e:
            logger.error(f"❌ 读取文本文件失败: {e}")
            raise

    def _process_pdf(self, file_path: str) -> str:
        """处理PDF文件 - 多方法 fallback（pdfplumber → PyPDF2 → PyMuPDF → OCR）"""
        if not self.pdf_available:
            raise Exception("PDF处理库未安装")

        content = ""

        # 方法1：pdfplumber（更好的表格支持）
        try:
            import pdfplumber
            logger.info("🔄 尝试使用pdfplumber提取PDF...")
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            content += f"[第{i + 1}页]\n{page_text}\n\n"
                    except Exception as e:
                        logger.warning(f"⚠️ 第{i + 1}页提取失败: {e}")

            if content.strip():
                logger.info(f"✅ pdfplumber成功提取 {len(content)} 字符")
                return self._clean_pdf_text(content)
        except ImportError:
            logger.warning("⚠️ pdfplumber未安装")
        except Exception as e:
            logger.warning(f"⚠️ pdfplumber处理失败: {e}")

        # 方法2：PyPDF2
        try:
            import PyPDF2
            logger.info("🔄 尝试使用PyPDF2提取PDF...")
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                num_pages = len(reader.pages)
                logger.info(f"📄 PDF共 {num_pages} 页")

                for page_num in range(num_pages):
                    try:
                        page = reader.pages[page_num]
                        page_text = page.extract_text()
                        if page_text:
                            content += f"[第{page_num + 1}页]\n{page_text}\n\n"
                    except Exception as e:
                        logger.warning(f"⚠️ 第{page_num + 1}页提取失败: {e}")

            if content.strip():
                logger.info(f"✅ PyPDF2成功提取 {len(content)} 字符")
                return self._clean_pdf_text(content)
        except ImportError:
            logger.warning("⚠️ PyPDF2未安装")
        except Exception as e:
            logger.warning(f"⚠️ PyPDF2处理失败: {e}")

        # 方法3：PyMuPDF (fitz)
        try:
            import fitz
            logger.info("🔄 尝试使用PyMuPDF提取PDF...")
            pdf_document = fitz.open(file_path)

            for page_num in range(pdf_document.page_count):
                try:
                    page = pdf_document[page_num]
                    page_text = page.get_text()
                    if page_text:
                        content += f"[第{page_num + 1}页]\n{page_text}\n\n"
                except Exception as e:
                    logger.warning(f"⚠️ 第{page_num + 1}页提取失败: {e}")

            pdf_document.close()

            if content.strip():
                logger.info(f"✅ PyMuPDF成功提取 {len(content)} 字符")
                return self._clean_pdf_text(content)
        except ImportError:
            logger.warning("⚠️ PyMuPDF未安装")
        except Exception as e:
            logger.warning(f"⚠️ PyMuPDF处理失败: {e}")

        # 方法4：OCR fallback
        if not content.strip():
            logger.warning("⚠️ 常规方法无法提取文本，尝试OCR...")
            content = self._try_ocr_pdf(file_path)
            if content:
                logger.info(f"✅ OCR成功提取 {len(content)} 字符")
                return content

        # 最终 fallback
        try:
            with open(file_path, 'rb') as f:
                header = f.read(4)
                if header == b'%PDF':
                    return f"[PDF文件: {Path(file_path).name}]\n[提取失败 - 可能是扫描版PDF]\n[建议转换为文本格式后重新上传]"
        except Exception:
            pass

        raise Exception("PDF文件无法提取文本内容（可能是扫描版或加密PDF）")

    def _clean_pdf_text(self, text: str) -> str:
        """清理PDF提取的文本"""
        lines = text.split('\n')
        cleaned_lines = []

        for line in lines:
            line = line.strip()
            if line:
                line = line.replace('\u200b', '').replace('\ufeff', '')
                line = ' '.join(line.split())
                cleaned_lines.append(line)

        result = '\n'.join(cleaned_lines)
        result = result.replace('？', '?')
        result = result.replace('！', '!')

        return result

    def _try_ocr_pdf(self, file_path: str) -> str:
        """尝试使用OCR识别PDF（需要额外依赖）"""
        try:
            import pytesseract
            from pdf2image import convert_from_path

            logger.info("🔄 开始OCR识别...")
            images = convert_from_path(file_path)

            content = []
            for i, image in enumerate(images):
                logger.info(f"🔍 OCR识别第 {i + 1}/{len(images)} 页...")
                text = pytesseract.image_to_string(image, lang='chi_sim+chi_tra+eng')
                if text.strip():
                    content.append(f"[第{i + 1}页]\n{text}")

            return '\n\n'.join(content)

        except ImportError:
            logger.warning("⚠️ OCR功能未安装。如需处理扫描版PDF，请安装 pytesseract 和 pdf2image。")
            return ""
        except Exception as e:
            logger.warning(f"⚠️ OCR识别失败: {e}")
            return ""

    def _process_docx(self, file_path: str) -> str:
        """处理DOCX文件 - 提取段落和表格"""
        if not self.docx_available:
            raise Exception("DOCX处理库未安装")

        try:
            import docx
            doc = docx.Document(file_path)

            content = []

            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content.append(paragraph.text)

            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        table_text.append(' | '.join(row_text))
                if table_text:
                    content.append('\n'.join(table_text))

            result = '\n\n'.join(content)
            logger.info(f"✅ 成功提取DOCX内容")
            return result

        except ImportError:
            raise Exception("请安装python-docx库: pip install python-docx")
        except Exception as e:
            logger.error(f"❌ DOCX处理失败: {e}")
            raise

    def _process_pptx(self, file_path: str) -> str:
        """处理PPTX文件"""
        if not self.pptx_available:
            raise Exception("PPTX处理库未安装")

        try:
            import pptx
            prs = pptx.Presentation(file_path)

            content = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = [f"=== 幻灯片 {slide_num} ==="]

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)

                if len(slide_content) > 1:
                    content.append('\n'.join(slide_content))

            result = '\n\n'.join(content)
            logger.info(f"✅ 成功提取PPTX内容")
            return result

        except ImportError:
            raise Exception("请安装python-pptx库: pip install python-pptx")
        except Exception as e:
            logger.error(f"❌ PPTX处理失败: {e}")
            raise

    def process_uploaded_file(self, file_data: bytes, filename: str,
                              user: str = "unknown") -> Tuple[bool, str, Dict]:
        """
        处理上传的文件数据（接受 bytes，内部创建临时文件）

        Args:
            file_data: 文件二进制数据
            filename: 原始文件名
            user: 上传用户

        Returns:
            (成功标志, 提取的文本内容, 文件信息字典)
        """
        temp_file = None
        try:
            ext = Path(filename).suffix.lower()

            with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
                tmp.write(file_data)
                temp_file = tmp.name

            return self.process_file(temp_file, filename, user)

        finally:
            if temp_file and os.path.exists(temp_file):
                try:
                    os.unlink(temp_file)
                except Exception:
                    pass


# 全局文件处理器实例（懒加载使用）
file_processor = FileProcessor()

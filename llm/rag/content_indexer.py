"""
学习中心内容索引器

负责将学习中心的各类内容（PDF、DOCX、PPTX、文章等）提取文本、
分割为 chunk 并存入 ChromaDB，供 RAG 检索使用。

架构：
    ContentTextExtractor — 纯文本提取，按 content_type 分派到对应解析器
    ContentIndexer       — 编排层：调用提取 → 分割 → 写入向量库

使用方式：
    indexer = ContentIndexer()
    indexer.index(content_id=42, content_row={...})
    indexer.is_indexed(content_id=42)
    indexer.delete(content_id=42)
"""

import logging
import os
from datetime import datetime
from typing import Dict, List, Optional, Tuple

from langchain.schema import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter

logger = logging.getLogger(__name__)

# 与 vector_store.py 保持一致的分割参数
CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200
SEPARATORS = ["\n\n", "\n", "。", "！", "？", "；", "，", " ", ""]

# ChromaDB metadata 中标识来源
SOURCE_TAG = "learning_center"


# ==================== 文本提取 ====================


class ContentTextExtractor:
    """
    从学习中心内容记录中提取纯文本。

    支持的格式：
        - article    → 直接使用 article_content 字段
        - document   → PDF (.pdf), DOCX (.docx), TXT/MD (.txt/.md)
        - image      → 不可提取文本，返回空串
        - video_*    → 不可提取文本，返回空串
    """

    def extract(self, content_row: Dict) -> str:
        """
        从内容记录提取纯文本。

        Args:
            content_row: lc_contents 表的完整行字典，需包含
                content_type, file_path, article_content, file_name, mime_type

        Returns:
            提取的文本。无法提取时返回空字符串。
        """
        content_type = content_row.get("content_type", "")

        if content_type == "article":
            return content_row.get("article_content") or ""

        if content_type in ("video_local", "video_external", "image"):
            return ""

        file_path = content_row.get("file_path", "")
        if not file_path or not os.path.exists(file_path):
            logger.warning("文件不存在或路径为空: %s", file_path)
            return ""

        return self._extract_from_file(file_path, content_row)

    def _extract_from_file(self, file_path: str, content_row: Dict) -> str:
        """根据文件类型分派到对应的提取方法。"""
        mime_type = content_row.get("mime_type") or ""
        file_name = (content_row.get("file_name") or file_path).lower()

        if mime_type == "application/pdf" or file_name.endswith(".pdf"):
            return self._extract_pdf(file_path)

        if (
            mime_type
            == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            or file_name.endswith(".docx")
        ):
            return self._extract_docx(file_path)

        if (
            mime_type
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            or file_name.endswith(".pptx")
        ):
            return self._extract_pptx(file_path)

        if file_name.endswith((".txt", ".md")):
            return self._extract_plaintext(file_path)

        logger.info(
            "不支持文本提取的文件类型: mime=%s, file=%s",
            mime_type,
            file_name,
        )
        return ""

    # ---- 按页提取（用于保留页码信息） ----

    def extract_with_pages(self, content_row: Dict) -> List[Dict]:
        """
        从内容记录提取按页分段的文本。

        返回 [{"page": 1, "text": "..."}, ...] 列表。
        非 PDF 类型或无"页"概念的内容，page 字段为 None。

        Args:
            content_row: lc_contents 表完整行字典

        Returns:
            按页分段的文本列表。无法提取时返回空列表。
        """
        content_type = content_row.get("content_type", "")

        if content_type == "article":
            text = content_row.get("article_content") or ""
            return [{"page": None, "text": text}] if text.strip() else []

        if content_type in ("video_local", "video_external", "image"):
            return []

        file_path = content_row.get("file_path", "")
        if not file_path or not os.path.exists(file_path):
            logger.warning("文件不存在或路径为空: %s", file_path)
            return []

        mime_type = content_row.get("mime_type") or ""
        file_name = (content_row.get("file_name") or file_path).lower()

        if mime_type == "application/pdf" or file_name.endswith(".pdf"):
            return self._extract_pdf_with_pages(file_path)

        # 非 PDF 文件（DOCX/PPTX/TXT）走原有逻辑，page=None
        text = self._extract_from_file(file_path, content_row)
        return [{"page": None, "text": text}] if text.strip() else []

    # ---- 各格式提取实现 ----

    @staticmethod
    def _extract_pdf(path: str) -> str:
        """提取 PDF 全文（不保留页码，向后兼容）。"""
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(path)
            pages = [page.get_text() for page in doc]
            doc.close()
            return "\n\n".join(pages)
        except Exception:
            logger.exception("PDF 文本提取失败: %s", path)
            return ""

    @staticmethod
    def _extract_pdf_with_pages(path: str) -> List[Dict]:
        """
        逐页提取 PDF 文本，保留页码信息。

        Returns:
            [{"page": 1, "text": "第一页文本"}, {"page": 2, "text": "..."}, ...]
            页码从 1 开始。空白页会被跳过。
        """
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(path)
            pages = []
            for page_num, page in enumerate(doc, 1):
                text = page.get_text().strip()
                if text:
                    pages.append({"page": page_num, "text": text})
            doc.close()
            logger.debug("PDF 逐页提取: %s, 有效页数=%d", path, len(pages))
            return pages
        except Exception:
            logger.exception("PDF 逐页提取失败: %s", path)
            return []

    @staticmethod
    def _extract_docx(path: str) -> str:
        try:
            from docx import Document as DocxDocument

            doc = DocxDocument(path)
            return "\n\n".join(
                p.text for p in doc.paragraphs if p.text.strip()
            )
        except Exception:
            logger.exception("DOCX 文本提取失败: %s", path)
            return ""

    @staticmethod
    def _extract_pptx(path: str) -> str:
        try:
            from pptx import Presentation

            prs = Presentation(path)
            slides: List[str] = []
            for slide_num, slide in enumerate(prs.slides, 1):
                texts = [
                    para.text.strip()
                    for shape in slide.shapes
                    if shape.has_text_frame
                    for para in shape.text_frame.paragraphs
                    if para.text.strip()
                ]
                if texts:
                    slides.append(
                        f"[Slide {slide_num}]\n" + "\n".join(texts)
                    )
            return "\n\n".join(slides)
        except Exception:
            logger.exception("PPTX 文本提取失败: %s", path)
            return ""

    @staticmethod
    def _extract_plaintext(path: str) -> str:
        try:
            with open(path, "r", encoding="utf-8") as f:
                return f.read()
        except Exception:
            logger.exception("文本文件读取失败: %s", path)
            return ""


# ==================== 内容索引器 ====================


class ContentIndexer:
    """
    学习中心内容的 ChromaDB 索引管理器。

    将内容文本分割为 chunk 后写入向量库，支持按 content_id 过滤检索。
    使用确定性 ID（lc_content_{id}_chunk_{i}），保证重复索引是幂等操作。
    """

    def __init__(self) -> None:
        self._extractor = ContentTextExtractor()
        self._splitter = RecursiveCharacterTextSplitter(
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP,
            separators=SEPARATORS,
        )

    def index(self, content_id: int, content_row: Dict) -> bool:
        """
        对一条学习内容建立向量索引。

        流程：文本提取（保留页码）→ chunk 分割 → 计算页码映射 → 写入 ChromaDB

        Args:
            content_id: lc_contents.id
            content_row: lc_contents 表完整行字典

        Returns:
            True 表示成功（包括内容无可提取文本的情况）。
            False 表示写入向量库时发生错误。
        """
        # 按页提取文本
        page_segments = self._extractor.extract_with_pages(content_row)
        if not page_segments:
            logger.info(
                "content_id=%s 无可索引文本（类型=%s），跳过",
                content_id,
                content_row.get("content_type", ""),
            )
            return True

        # 拼接全文并构建字符偏移 → 页码映射
        full_text, offset_to_page = self._build_text_with_page_map(page_segments)
        if not full_text.strip():
            return True

        chunks = self._splitter.split_text(full_text)
        if not chunks:
            logger.warning("content_id=%s 文本分割后无 chunk", content_id)
            return True

        doc_id = f"lc_content_{content_id}"
        title = content_row.get("title", "")
        filename = content_row.get("file_name", "")

        # 为每个 chunk 计算覆盖的页码
        chunk_page_map = self._map_chunks_to_pages(
            full_text, chunks, offset_to_page
        )

        documents = [
            Document(
                page_content=chunk,
                metadata={
                    "doc_id": doc_id,
                    "chunk_id": i,
                    "content_id": str(content_id),
                    "title": title,
                    "filename": filename,
                    "source": SOURCE_TAG,
                    "total_chunks": len(chunks),
                    "indexed_at": datetime.now().isoformat(),
                    "page_numbers": chunk_page_map.get(i, ""),
                },
            )
            for i, chunk in enumerate(chunks)
        ]
        ids = [f"{doc_id}_chunk_{i}" for i in range(len(documents))]

        try:
            vector_db = _get_vector_db()
            vector_db.add_documents(documents, ids=ids)
            logger.info(
                "已索引 content_id=%s: %d chunks（标题='%s'）",
                content_id,
                len(chunks),
                title,
            )
            return True
        except Exception:
            logger.exception("content_id=%s 向量库写入失败", content_id)
            return False

    @staticmethod
    def _build_text_with_page_map(
        page_segments: List[Dict],
    ) -> Tuple[str, List[Tuple[int, int, Optional[int]]]]:
        """
        将按页分段的文本拼接为全文，并生成字符偏移→页码映射。

        Returns:
            (full_text, offset_map)
            offset_map: [(start, end, page_number), ...] 按 start 升序
            page_number 为 None 表示该段无页码信息（非 PDF）
        """
        parts: List[str] = []
        offset_map: List[Tuple[int, int, Optional[int]]] = []
        cursor = 0

        for seg in page_segments:
            text = seg["text"]
            page = seg.get("page")
            if not text.strip():
                continue

            start = cursor
            parts.append(text)
            cursor += len(text)
            offset_map.append((start, cursor, page))

            # 段间分隔符
            parts.append("\n\n")
            cursor += 2

        full_text = "".join(parts).rstrip("\n")
        return full_text, offset_map

    @staticmethod
    def _map_chunks_to_pages(
        full_text: str,
        chunks: List[str],
        offset_map: List[Tuple[int, int, Optional[int]]],
    ) -> Dict[int, str]:
        """
        根据 chunk 在全文中的位置，映射到覆盖的页码范围。

        Returns:
            {chunk_index: "3,4"} — 逗号分隔的页码字符串。
            无页码信息的 chunk 对应空字符串。
        """
        result: Dict[int, str] = {}
        search_start = 0

        for i, chunk in enumerate(chunks):
            # 在全文中查找 chunk 的起始位置
            pos = full_text.find(chunk, search_start)
            if pos == -1:
                # 回退到全文搜索（处理 overlap 导致的位置偏移）
                pos = full_text.find(chunk)
            if pos == -1:
                result[i] = ""
                continue

            chunk_start = pos
            chunk_end = pos + len(chunk)
            search_start = pos  # 下次从当前位置开始找

            # 找出与 chunk 区间重叠的所有页码
            pages = set()
            for seg_start, seg_end, page in offset_map:
                if page is not None and seg_start < chunk_end and seg_end > chunk_start:
                    pages.add(page)

            result[i] = ",".join(str(p) for p in sorted(pages))

        return result

    def is_indexed(self, content_id: int) -> bool:
        """检查某条内容是否已有向量索引。"""
        try:
            collection = _get_vector_db()._collection
            results = collection.get(
                where={"content_id": str(content_id)},
                limit=1,
            )
            return len(results.get("ids", [])) > 0
        except Exception:
            logger.exception(
                "检查索引状态失败: content_id=%s", content_id
            )
            return False

    def has_page_metadata(self, content_id: int) -> bool:
        """
        检查已索引的内容是否包含 page_numbers metadata。

        用于判断旧索引是否需要重建（v1.8.0 之前的索引缺少此字段）。
        """
        try:
            collection = _get_vector_db()._collection
            results = collection.get(
                where={"content_id": str(content_id)},
                limit=1,
                include=["metadatas"],
            )
            metadatas = results.get("metadatas", [])
            if not metadatas:
                return False
            return "page_numbers" in metadatas[0]
        except Exception:
            logger.exception(
                "检查页码元数据失败: content_id=%s", content_id
            )
            return False

    def delete(self, content_id: int) -> bool:
        """删除某条内容的所有向量索引。"""
        try:
            collection = _get_vector_db()._collection
            collection.delete(where={"content_id": str(content_id)})
            logger.info("已删除 content_id=%s 的向量索引", content_id)
            return True
        except Exception:
            logger.exception(
                "删除索引失败: content_id=%s", content_id
            )
            return False


# ==================== 模块级便捷函数 ====================

# 模块级单例，避免每次调用都创建 splitter
_indexer: Optional[ContentIndexer] = None


def get_content_indexer() -> ContentIndexer:
    """获取 ContentIndexer 单例。"""
    global _indexer
    if _indexer is None:
        _indexer = ContentIndexer()
    return _indexer


def _get_vector_db():
    """获取向量库实例（委托给 retrieval 模块的单例）。"""
    from .retrieval import get_vector_db

    return get_vector_db()
